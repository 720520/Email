"""生成内置竖版周报，或把数据填入用户上传的 PPTX 模板。"""

from __future__ import annotations

import base64
import io
import math
from collections.abc import Iterable
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.services.report_template_service import (
    COMPONENT_PATTERN,
    TOKEN_PATTERN,
    ReportTemplateInspector,
    format_template_value,
)

_BLUE = RGBColor(64, 113, 190)
_ORANGE = RGBColor(247, 130, 35)
_DARK = RGBColor(55, 55, 55)
_LIGHT_ORANGE = RGBColor(253, 242, 234)
_BORDER = RGBColor(190, 190, 190)

_PRODUCT_HEADER_FIELDS = {
    "产品名称": "product_name",
    "成立日期": "inception_date",
    "策略分类": "strategy_category",
    "年化收益率": "annualized_return",
    "夏普比率": "sharpe_ratio",
    "最新单位净值": "latest_unit_nav",
    "最新累计净值": "latest_total_nav",
}
_PERFORMANCE_HEADER_FIELDS = {
    "近一月": "return_1m",
    "近三月": "return_3m",
    "近半年": "return_6m",
    "今年以来": "return_ytd",
    "近一年": "return_1y",
    "成立以来": "return_since",
    "今年以来最大回撤": "max_drawdown_ytd",
    "成立以来最大回撤": "max_drawdown_since",
}
_CONTRACT_LABEL_FIELDS = {
    "管理机构": "manager_name",
    "托管/外包机构": "custodian_name",
    "风险登记": "risk_level",
    "风险等级": "risk_level",
    "开放日": "open_day",
    "存续期间": "duration",
    "存续期限": "duration",
    "锁定期": "lockup_period",
    "业绩报酬": "performance_fee",
    "投资范围": "investment_scope",
    "本基金承担费率": "fund_fees",
    "投资者承担费率": "investor_fees",
}


class ReportPresentationService:
    def generate(
        self,
        snapshot: dict[str, Any],
        *,
        output_path: Path,
        sections: list[str],
        template_path: Path | None = None,
    ) -> None:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        if template_path is None:
            presentation = self._build_standard(snapshot, sections)
        else:
            presentation = Presentation(str(template_path))
            self._fill_template(presentation, snapshot)
        presentation.save(str(output_path))

    def _build_standard(self, snapshot: dict[str, Any], sections: list[str]) -> Presentation:
        presentation = Presentation()
        presentation.slide_width = Inches(8.27)
        presentation.slide_height = Inches(11.69)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        fields = snapshot["fields"]
        metrics = snapshot["performance"]

        header = slide.shapes.add_shape(1, 0, 0, presentation.slide_width, Inches(0.58))
        header.fill.solid()
        header.fill.fore_color.rgb = _BLUE
        header.line.fill.background()
        self._add_text(
            slide,
            fields.get("product_name") or snapshot["product_name"],
            0.82,
            0.06,
            5.8,
            0.42,
            size=23,
            bold=True,
            color=RGBColor(255, 255, 255),
        )
        self._add_text(
            slide,
            f"基金周报  {snapshot['report_date']}",
            6.2,
            0.12,
            1.85,
            0.28,
            size=10,
            bold=True,
            color=RGBColor(255, 255, 255),
            align=PP_ALIGN.RIGHT,
        )
        y = 0.72

        if "product_info" in sections:
            y = self._section_title(slide, "产品信息", y)
            headers = list(_PRODUCT_HEADER_FIELDS)
            values = [
                self._value_for(field_key, fields, metrics)
                for field_key in _PRODUCT_HEADER_FIELDS.values()
            ]
            self._add_table(slide, headers, [values], 0.2, y, 7.87, 0.58, font_size=8.5)
            y += 0.72

        if "performance" in sections:
            y = self._section_title(slide, "收益指标（周频）", y)
            headers = ["收益指标", *_PERFORMANCE_HEADER_FIELDS.keys()]
            values = [
                snapshot["product_name"],
                *[
                    metrics.get(field_key) or "—"
                    for field_key in _PERFORMANCE_HEADER_FIELDS.values()
                ],
            ]
            self._add_table(slide, headers, [values], 0.2, y, 7.87, 0.67, font_size=7.7)
            y += 0.81

        if "nav_chart" in sections:
            y = self._section_title(slide, "业绩曲线", y)
            chart_height = 3.25 if "strategy" in sections else 4.2
            self._add_nav_chart(slide, snapshot, 0.24, y, 7.78, chart_height)
            y += chart_height + 0.12

        if "strategy" in sections:
            y = self._section_title(slide, "策略介绍", y)
            strategy = fields.get("investment_strategy") or "尚未维护策略介绍"
            self._add_text(slide, strategy, 0.26, y, 7.7, 1.0, size=9.3, color=_DARK)
            y += 1.12

        if "contract_terms" in sections:
            rows = [
                [
                    "管理机构",
                    fields.get("manager_name") or "—",
                    "托管/外包机构",
                    fields.get("custodian_name") or "—",
                ],
                [
                    "风险等级",
                    fields.get("risk_level") or "—",
                    "开放日",
                    fields.get("open_day") or "—",
                ],
                [
                    "存续期间",
                    fields.get("duration") or "—",
                    "锁定期",
                    fields.get("lockup_period") or "—",
                ],
                [
                    "本基金承担费率",
                    self._fund_fees(fields),
                    "投资者承担费率",
                    self._investor_fees(fields),
                ],
                [
                    "业绩报酬",
                    fields.get("performance_fee") or "—",
                    "投资范围",
                    fields.get("investment_scope") or "—",
                ],
            ]
            available = max(1.55, min(2.15, 11.15 - y))
            self._add_matrix_table(slide, rows, 0.13, y, 8.0, available)
            y += available + 0.08

        if "disclaimer" in sections:
            disclaimer = fields.get("disclaimer") or (
                "声明：本材料仅供内部及特定客户参阅，不构成投资建议。过往业绩并不预示未来表现，"
                "产品信息以基金合同及最新法律文件为准。"
            )
            self._add_text(slide, disclaimer, 0.18, min(y, 11.2), 7.9, 0.35, size=6.8, color=_DARK)
        return presentation

    def _fill_template(self, presentation: Presentation, snapshot: dict[str, Any]) -> None:
        fields = snapshot["fields"]
        metrics = snapshot["performance"]
        tokens = {
            **{key: value or "—" for key, value in fields.items()},
            **{key: value or "—" for key, value in metrics.items()},
            "report_date": snapshot["report_date"],
            "fund_fees": self._fund_fees(fields),
            "investor_fees": self._investor_fees(fields),
            "product.name": snapshot["product_name"],
            "product.code": snapshot.get("product_code") or fields.get("product_code"),
            "product.investment_manager": fields.get("investment_manager"),
            "product.investment_strategy": fields.get("investment_strategy"),
            "product.inception_date": fields.get("inception_date"),
            "product.manager_name": fields.get("manager_name"),
            "product.custodian_name": fields.get("custodian_name"),
            "report.date": snapshot["report_date"],
            **{f"metric.{key}": value for key, value in metrics.items()},
            **(snapshot.get("dynamic_fields") or {}),
        }
        for slide in presentation.slides:
            strategy_heading_bottom = None
            contract_table_top = None
            for shape in list(self._all_shapes(slide.shapes)):
                if getattr(shape, "has_text_frame", False):
                    original = shape.text.strip()
                    if self._render_component_anchor(slide, shape, original, snapshot, tokens):
                        continue
                    self._replace_tokens(shape.text_frame, tokens)
                    if original == "策略介绍":
                        strategy_heading_bottom = shape.top + shape.height
                    if original.startswith("产品报告日期"):
                        report_date = snapshot["report_date"].replace("-", "/")
                        self._set_text(shape, f"产品报告日期：{report_date}")
                    elif (
                        shape.top < Inches(0.8)
                        and "私募证券投资基金" in original
                        and not original.startswith("声明")
                    ):
                        self._set_text(shape, snapshot["product_name"])
                    elif (
                        (fields.get("slogan") or fields.get("strategy_category"))
                        and 4 <= len(original) <= 60
                        and "CTA" in original.upper()
                    ):
                        summary = fields.get("slogan") or fields.get("strategy_category")
                        self._set_text(shape, str(summary))
                    elif original.startswith("声明："):
                        disclaimer = fields.get("disclaimer") or (
                            "声明：本材料仅供内部及特定客户参阅，不构成投资建议。过往业绩"
                            "并不预示未来表现，产品信息以基金合同及最新法律文件为准。"
                        )
                        self._set_text(shape, str(disclaimer))
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            self._replace_tokens(cell.text_frame, tokens)
                    table_text = "|".join(
                        cell.text for row in shape.table.rows for cell in row.cells
                    )
                    if "管理机构" in table_text or "本基金承担费率" in table_text:
                        contract_table_top = shape.top
                    self._fill_table(shape.table, snapshot)
                if getattr(shape, "has_chart", False):
                    self._replace_chart(slide, shape, snapshot)
            if strategy_heading_bottom is not None and fields.get("investment_strategy"):
                candidates = [
                    shape
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                    and shape.top >= strategy_heading_bottom
                    and (contract_table_top is None or shape.top < contract_table_top)
                    and len(shape.text.strip()) >= 60
                    and not shape.text.strip().startswith("声明")
                ]
                if candidates:
                    target = max(candidates, key=lambda item: len(item.text))
                    self._set_text(target, str(fields["investment_strategy"]))

    def _fill_table(self, table, snapshot: dict[str, Any]) -> None:
        fields = snapshot["fields"]
        metrics = snapshot["performance"]
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        if "产品名称" in headers and len(table.rows) >= 2:
            for column, header in enumerate(headers):
                field_key = _PRODUCT_HEADER_FIELDS.get(header)
                if field_key:
                    self._set_cell_text(
                        table.cell(1, column),
                        str(self._value_for(field_key, fields, metrics)),
                    )
        elif "收益指标" in headers and len(table.rows) >= 2:
            self._set_cell_text(table.cell(1, 0), snapshot["product_name"])
            for column, header in enumerate(headers[1:], start=1):
                field_key = _PERFORMANCE_HEADER_FIELDS.get(header)
                if field_key:
                    self._set_cell_text(table.cell(1, column), str(metrics.get(field_key) or "—"))
            # 模板中的旧基准数据没有当前系统来源时必须清空，禁止串用历史产品数据。
            for row_index in range(2, len(table.rows)):
                for cell in table.rows[row_index].cells:
                    self._set_cell_text(cell, "")
        else:
            if table.rows and "私募证券投资基金" in table.cell(0, 0).text:
                self._set_cell_text(table.cell(0, 0), snapshot["product_name"])
            for row in table.rows:
                for label_column in range(0, len(row.cells) - 1, 2):
                    field_key = _CONTRACT_LABEL_FIELDS.get(row.cells[label_column].text.strip())
                    if not field_key:
                        continue
                    value = (
                        self._fund_fees(fields)
                        if field_key == "fund_fees"
                        else self._investor_fees(fields)
                        if field_key == "investor_fees"
                        else fields.get(field_key) or "—"
                    )
                    self._set_cell_text(row.cells[label_column + 1], str(value))

    @staticmethod
    def _replace_chart(slide, shape, snapshot: dict[str, Any]) -> None:
        series = snapshot["nav_series"]
        if not series:
            return
        values = [float(item["total_nav"] or item["unit_nav"]) for item in series]
        chart_data = ChartData()
        chart_data.categories = [date.fromisoformat(item["date"]) for item in series]
        chart_data.add_series(
            f"{snapshot['product_name']} 累计净值",
            values,
        )
        try:
            shape.chart.replace_data(chart_data)
            ReportPresentationService._configure_value_axis(shape.chart, values)
            shape.chart.category_axis.tick_labels.number_format = "yyyy-mm-dd"
            shape.chart.category_axis.tick_label_skip = max(1, len(series) // 10)
        except ValueError:
            # 部分托管模板图表引用外部 Excel，python-pptx 无法修改其外链工作簿。
            # 此时保留模板中的图表区域，用本地内嵌数据图表原位替换。
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            shape.element.getparent().remove(shape.element)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, left, top, width, height, chart_data
            ).chart
            chart.has_title = True
            chart.chart_title.text_frame.text = f"{snapshot['product_name']} 累计净值"
            chart.has_legend = False
            chart.value_axis.has_major_gridlines = True
            ReportPresentationService._configure_value_axis(chart, values)
            chart.category_axis.tick_labels.number_format = "yyyy-mm-dd"
            chart.category_axis.tick_label_skip = max(1, len(series) // 10)
            chart.category_axis.tick_labels.font.size = Pt(6)
            chart.value_axis.tick_labels.font.size = Pt(7)
            chart.series[0].format.line.color.rgb = _DARK
            chart.series[0].format.line.width = Pt(1.8)

    def _add_nav_chart(self, slide, snapshot: dict[str, Any], x, y, width, height) -> None:
        chart_data = ChartData()
        series = snapshot["nav_series"]
        values = [float(item["total_nav"] or item["unit_nav"]) for item in series]
        chart_data.categories = [date.fromisoformat(item["date"]) for item in series]
        chart_data.add_series("累计净值", values)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(width), Inches(height), chart_data
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = f"{snapshot['product_name']} 累计净值"
        chart.has_legend = False
        chart.value_axis.has_major_gridlines = True
        self._configure_value_axis(chart, values)
        chart.category_axis.tick_labels.number_format = "yyyy-mm-dd"
        chart.category_axis.tick_label_skip = max(1, len(series) // 10)
        chart.category_axis.tick_labels.font.size = Pt(6)
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.series[0].format.line.color.rgb = _DARK
        chart.series[0].format.line.width = Pt(1.8)

    def _section_title(self, slide, title: str, y: float) -> float:
        self._add_text(slide, title, 0.2, y, 3.4, 0.3, size=15, bold=True, color=_BLUE)
        line = slide.shapes.add_shape(1, Inches(0.2), Inches(y + 0.34), Inches(7.87), Pt(1.2))
        line.fill.solid()
        line.fill.fore_color.rgb = _ORANGE
        line.line.fill.background()
        return y + 0.43

    def _add_table(
        self,
        slide,
        headers: list[str],
        rows: list[list[str]],
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        font_size: float,
    ) -> None:
        table = slide.shapes.add_table(
            len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(width), Inches(height)
        ).table
        for column, header in enumerate(headers):
            table.cell(0, column).text = header
        for row_index, values in enumerate(rows, start=1):
            for column, value in enumerate(values):
                table.cell(row_index, column).text = str(value or "—")
        for row in table.rows:
            for cell in row.cells:
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = "微软雅黑"
                        run.font.size = Pt(font_size)

    def _add_matrix_table(
        self, slide, rows: list[list[str]], x: float, y: float, width: float, height: float
    ) -> None:
        table = slide.shapes.add_table(
            len(rows), 4, Inches(x), Inches(y), Inches(width), Inches(height)
        ).table
        widths = [1.35, 2.25, 1.35, 3.05]
        for index, width_value in enumerate(widths):
            table.columns[index].width = Inches(width_value)
        for row_index, values in enumerate(rows):
            for column, value in enumerate(values):
                cell = table.cell(row_index, column)
                cell.text = str(value or "—")
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    _LIGHT_ORANGE if row_index % 2 else RGBColor(255, 255, 255)
                )
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = "微软雅黑"
                        run.font.size = Pt(7.8)
        table.cell(0, 0).margin_left = Pt(2)

    @staticmethod
    def _add_text(
        slide,
        text: str,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        size: float,
        bold: bool = False,
        color: RGBColor = _DARK,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(text or "—")
        run.font.name = "微软雅黑"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    @staticmethod
    def _replace_tokens(text_frame, tokens: dict[str, Any]) -> None:
        for paragraph in text_frame.paragraphs:
            original = paragraph.text
            if "{{" not in original:
                continue

            def replace(match) -> str:
                expression = match.group(1).strip()
                parts = [part.strip() for part in expression.split("|")]
                target = parts[0]
                if COMPONENT_PATTERN.match(target):
                    return match.group(0)
                normalized = ReportTemplateInspector.normalize_field_key(target)
                value = tokens.get(normalized, tokens.get(target))
                return format_template_value(value, tuple(parts[1:]))

            rendered = TOKEN_PATTERN.sub(replace, original)
            ReportPresentationService._set_paragraph_text(paragraph, rendered)

    def _render_component_anchor(
        self,
        slide,
        shape,
        text: str,
        snapshot: dict[str, Any],
        tokens: dict[str, Any],
    ) -> bool:
        match = TOKEN_PATTERN.fullmatch(text)
        if not match:
            return False
        target = match.group(1).strip().split("|", 1)[0].strip()
        component_match = COMPONENT_PATTERN.match(target)
        if not component_match:
            return False
        kind, name = component_match.groups()
        x = shape.left / Inches(1)
        y = shape.top / Inches(1)
        width = shape.width / Inches(1)
        height = shape.height / Inches(1)
        shape.element.getparent().remove(shape.element)
        if kind == "chart" and name == "nav_history":
            self._add_nav_chart(slide, snapshot, x, y, width, height)
        elif kind == "table" and name == "product_info":
            fields, metrics = snapshot["fields"], snapshot["performance"]
            headers = list(_PRODUCT_HEADER_FIELDS)
            values = [
                self._value_for(field_key, fields, metrics)
                for field_key in _PRODUCT_HEADER_FIELDS.values()
            ]
            self._add_table(slide, headers, [values], x, y, width, height, font_size=8.5)
        elif kind == "table" and name == "performance":
            metrics = snapshot["performance"]
            headers = ["收益指标", *_PERFORMANCE_HEADER_FIELDS.keys()]
            values = [
                snapshot["product_name"],
                *[metrics.get(key) or "—" for key in _PERFORMANCE_HEADER_FIELDS.values()],
            ]
            self._add_table(slide, headers, [values], x, y, width, height, font_size=7.7)
        elif kind == "image":
            field_key = ReportTemplateInspector.normalize_field_key(name.strip())
            self._add_data_image(slide, tokens.get(field_key), x, y, width, height)
        return True

    @staticmethod
    def _add_data_image(slide, value: Any, x: float, y: float, width: float, height: float) -> None:
        if (
            not isinstance(value, str)
            or not value.startswith("data:image/")
            or ";base64," not in value
        ):
            ReportPresentationService._add_text(
                slide, "图片未配置", x, y, width, height, size=9, color=_DARK
            )
            return
        try:
            content = base64.b64decode(value.split(";base64,", 1)[1], validate=True)
        except ValueError:
            content = b""
        if not content or len(content) > 20 * 1024 * 1024:
            ReportPresentationService._add_text(
                slide, "图片数据无效", x, y, width, height, size=9, color=_DARK
            )
            return
        slide.shapes.add_picture(
            io.BytesIO(content), Inches(x), Inches(y), Inches(width), Inches(height)
        )

    @staticmethod
    def _set_paragraph_text(paragraph, value: str) -> None:
        if paragraph.runs:
            paragraph.runs[0].text = value
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = value

    @staticmethod
    def _set_cell_text(cell, value: str) -> None:
        ReportPresentationService._set_text_frame_text(cell.text_frame, value)

    @staticmethod
    def _set_text(shape, value: str) -> None:
        ReportPresentationService._set_text_frame_text(shape.text_frame, value)

    @staticmethod
    def _set_text_frame_text(frame, value: str) -> None:
        if frame.paragraphs and frame.paragraphs[0].runs:
            frame.paragraphs[0].runs[0].text = value
            for run in frame.paragraphs[0].runs[1:]:
                run.text = ""
            for paragraph in frame.paragraphs[1:]:
                for run in paragraph.runs:
                    run.text = ""
        else:
            frame.text = value

    @staticmethod
    def _configure_value_axis(chart, values: list[float]) -> None:
        span = max(values) - min(values)
        padding = span * 0.08 if span else max(abs(values[0]) * 0.02, 0.01)
        minimum = max(0, math.floor((min(values) - padding) * 20) / 20)
        maximum = math.ceil((max(values) + padding) * 20) / 20
        chart.value_axis.minimum_scale = minimum
        chart.value_axis.maximum_scale = maximum
        chart.value_axis.major_unit = max((maximum - minimum) / 5, 0.01)
        chart.value_axis.tick_labels.number_format = "0.0000"

    @staticmethod
    def _all_shapes(shapes) -> Iterable:
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from ReportPresentationService._all_shapes(shape.shapes)

    @staticmethod
    def _fund_fees(fields: dict[str, Any]) -> str:
        values = [
            f"管理费 {fields['management_fee']}" if fields.get("management_fee") else None,
            f"托管/外包费 {fields['custody_fee']}" if fields.get("custody_fee") else None,
        ]
        return "\n".join(item for item in values if item) or "—"

    @staticmethod
    def _investor_fees(fields: dict[str, Any]) -> str:
        values = [
            f"申购费 {fields['subscription_fee']}" if fields.get("subscription_fee") else None,
            f"赎回费 {fields['redemption_fee']}" if fields.get("redemption_fee") else None,
        ]
        return "\n".join(item for item in values if item) or "—"

    @staticmethod
    def _value_for(field_key: str, fields: dict[str, Any], metrics: dict[str, Any]) -> str:
        if field_key in metrics:
            return str(metrics.get(field_key) or "—")
        return str(fields.get(field_key) or "—")
