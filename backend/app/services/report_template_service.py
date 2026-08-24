"""PPTX 模板占位符发现、字段依赖校验与格式化。"""

from __future__ import annotations

import re
from collections.abc import Iterable
from dataclasses import dataclass
from datetime import date, datetime
from decimal import Decimal, InvalidOperation
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from sqlalchemy.orm import Session

from app.core.errors import AppError
from app.services.report_field_service import ReportFieldResolver
from app.services.reporting_service import REPORT_FIELD_DEFINITIONS

TOKEN_PATTERN = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")
FIELD_KEY_PATTERN = re.compile(r"^[a-z][a-z0-9_]*(?:\.[a-z][a-z0-9_]*)*$")
COMPONENT_PATTERN = re.compile(r"^(table|chart|image):(.+)$")
PERCENT_FORMAT = re.compile(r"^percent(?::([0-6]))?$")
DEFAULT_FORMAT = re.compile(r'^default:"([^"]*)"$')

LEGACY_FIELD_ALIASES: dict[str, str] = {
    "product_name": "product.name",
    "product_code": "product.code",
    "investment_manager": "product.investment_manager",
    "investment_strategy": "product.investment_strategy",
    "inception_date": "product.inception_date",
    "manager_name": "product.manager_name",
    "custodian_name": "product.custodian_name",
    "report_date": "report.date",
    "annualized_return": "metric.annualized_return",
    "sharpe_ratio": "metric.sharpe_ratio",
    "return_ytd": "metric.return_ytd",
    "max_drawdown_since": "metric.max_drawdown_since",
    "company_logo": "custom.company_logo",
}
LEGACY_SNAPSHOT_FIELDS = set(REPORT_FIELD_DEFINITIONS) | {
    "latest_unit_nav",
    "latest_total_nav",
    "return_1m",
    "return_3m",
    "return_6m",
    "return_1y",
    "return_since",
    "max_drawdown_ytd",
    "fund_fees",
    "investor_fees",
}
KNOWN_COMPONENTS = {"table:product_info", "table:performance", "chart:nav_history"}
LEGACY_FIELD_TYPES = {
    "report_date": "date",
    "annualized_return": "percentage",
    "return_ytd": "percentage",
    "max_drawdown_since": "percentage",
}


@dataclass(frozen=True, slots=True)
class TemplateToken:
    expression: str
    field_key: str | None
    component: str | None
    formatters: tuple[str, ...]
    slide_number: int
    location: str


@dataclass(frozen=True, slots=True)
class TemplateInspection:
    required_fields: tuple[str, ...]
    required_components: tuple[str, ...]
    tokens: tuple[TemplateToken, ...]
    errors: tuple[dict[str, Any], ...]

    @property
    def is_valid(self) -> bool:
        return not self.errors


class ReportTemplateInspector:
    def inspect_path(self, path: Path, session: Session | None = None) -> TemplateInspection:
        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise AppError("REPORT_TEMPLATE_INVALID", "PPTX 模板无法打开或结构已损坏") from exc
        return self.inspect(presentation, session=session)

    def inspect(
        self, presentation: Presentation, session: Session | None = None
    ) -> TemplateInspection:
        tokens: list[TemplateToken] = []
        errors: list[dict[str, Any]] = []
        component_counts: dict[str, int] = {}
        scopes: list[tuple[int, str, Any]] = [
            (slide_number, "slide", slide.shapes)
            for slide_number, slide in enumerate(presentation.slides, start=1)
        ]
        scopes.extend(
            (0, f"master:{master_index}", master.shapes)
            for master_index, master in enumerate(presentation.slide_masters, start=1)
        )
        scopes.extend(
            (0, f"layout:{layout_index}", layout.shapes)
            for layout_index, layout in enumerate(presentation.slide_layouts, start=1)
        )
        for slide_number, scope_name, shapes in scopes:
            for location, text in self._text_blocks(shapes, prefix=scope_name):
                matches = list(TOKEN_PATTERN.finditer(text))
                scrubbed = TOKEN_PATTERN.sub("", text)
                if "{{" in scrubbed or "}}" in scrubbed:
                    errors.append(
                        self._error("MALFORMED_TOKEN", slide_number, location, "占位符大括号不完整")
                    )
                for match in matches:
                    token, token_errors = self._parse_token(
                        match.group(1).strip(), slide_number, location, session
                    )
                    errors.extend(token_errors)
                    if token:
                        tokens.append(token)
                        if token.component:
                            component_counts[token.component] = (
                                component_counts.get(token.component, 0) + 1
                            )
        for component, count in component_counts.items():
            if count > 1:
                errors.append(
                    {
                        "code": "DUPLICATE_COMPONENT",
                        "message": f"结构化组件 {component} 重复出现 {count} 次",
                        "component": component,
                    }
                )
        return TemplateInspection(
            required_fields=tuple(sorted({item.field_key for item in tokens if item.field_key})),
            required_components=tuple(sorted(component_counts)),
            tokens=tuple(tokens),
            errors=tuple(errors),
        )

    def _parse_token(
        self, expression: str, slide: int, location: str, session: Session | None
    ) -> tuple[TemplateToken | None, list[dict[str, Any]]]:
        parts = [part.strip() for part in expression.split("|")]
        target, formatters = parts[0], tuple(parts[1:])
        errors: list[dict[str, Any]] = []
        component_match = COMPONENT_PATTERN.match(target)
        if component_match:
            kind, name = component_match.groups()
            component = f"{kind}:{name.strip()}"
            if kind == "image":
                field_key = self.normalize_field_key(name.strip())
                definition, field_errors = self._field_definition(
                    field_key, slide, location, session
                )
                errors.extend(field_errors)
                if definition and definition["data_type"] != "image":
                    errors.append(
                        self._error(
                            "IMAGE_FIELD_TYPE_MISMATCH",
                            slide,
                            location,
                            f"图片锚点字段必须是 image 类型：{field_key}",
                        )
                    )
            elif component not in KNOWN_COMPONENTS:
                errors.append(
                    self._error(
                        "UNKNOWN_COMPONENT", slide, location, f"未知结构化组件：{component}"
                    )
                )
                field_key = None
            else:
                field_key = None
            if formatters:
                errors.append(
                    self._error(
                        "COMPONENT_FORMAT_UNSUPPORTED", slide, location, "结构化组件不支持格式化器"
                    )
                )
            return TemplateToken(expression, field_key, component, (), slide, location), errors

        field_key = self.normalize_field_key(target)
        if not FIELD_KEY_PATTERN.fullmatch(field_key):
            errors.append(
                self._error("INVALID_FIELD_KEY", slide, location, f"字段标识格式无效：{target}")
            )
        else:
            definition, field_errors = self._field_definition(field_key, slide, location, session)
            errors.extend(field_errors)
            errors.extend(
                self._validate_formatters(field_key, definition, formatters, slide, location)
            )
        for formatter in formatters:
            if (
                formatter == "date"
                or PERCENT_FORMAT.fullmatch(formatter)
                or DEFAULT_FORMAT.fullmatch(formatter)
            ):
                continue
            errors.append(
                self._error("UNKNOWN_FORMATTER", slide, location, f"未知格式化器：{formatter}")
            )
        return TemplateToken(expression, field_key, None, formatters, slide, location), errors

    @staticmethod
    def normalize_field_key(key: str) -> str:
        return LEGACY_FIELD_ALIASES.get(key, key)

    @staticmethod
    def _field_definition(
        field_key: str, slide: int, location: str, session: Session | None
    ) -> tuple[dict[str, Any] | None, list[dict[str, Any]]]:
        if field_key in LEGACY_SNAPSHOT_FIELDS:
            return {"data_type": LEGACY_FIELD_TYPES.get(field_key, "string")}, []
        if session is None:
            return None, []
        try:
            definition = ReportFieldResolver().definition(session, field_key)
        except AppError:
            return None, [
                {
                    "code": "UNKNOWN_FIELD",
                    "slide": slide,
                    "location": location,
                    "field_key": field_key,
                    "message": f"字段未注册或已停用：{field_key}",
                }
            ]
        errors = []
        if definition.get("is_required") and definition.get("default_value") is None:
            errors.append(
                {
                    "code": "REQUIRED_FIELD_DEFAULT_MISSING",
                    "slide": slide,
                    "location": location,
                    "field_key": field_key,
                    "message": f"必填字段未配置默认值：{field_key}",
                }
            )
        return definition, errors

    @classmethod
    def _validate_formatters(
        cls,
        field_key: str,
        definition: dict[str, Any] | None,
        formatters: tuple[str, ...],
        slide: int,
        location: str,
    ) -> list[dict[str, Any]]:
        if definition is None:
            return []
        data_type = definition["data_type"]
        errors = []
        for formatter in formatters:
            if formatter == "date" and data_type != "date":
                errors.append(
                    cls._error(
                        "FORMATTER_TYPE_MISMATCH",
                        slide,
                        location,
                        f"date 格式化器不能用于 {data_type} 字段：{field_key}",
                    )
                )
            if PERCENT_FORMAT.fullmatch(formatter) and data_type not in {
                "number",
                "percentage",
            }:
                errors.append(
                    cls._error(
                        "FORMATTER_TYPE_MISMATCH",
                        slide,
                        location,
                        f"percent 格式化器不能用于 {data_type} 字段：{field_key}",
                    )
                )
        return errors

    @staticmethod
    def _error(code: str, slide: int, location: str, message: str) -> dict[str, Any]:
        return {"code": code, "slide": slide, "location": location, "message": message}

    def _text_blocks(self, shapes, prefix: str = "") -> Iterable[tuple[str, str]]:
        for shape_index, shape in enumerate(shapes, start=1):
            location = f"{prefix}/shape:{shape_index}" if prefix else f"shape:{shape_index}"
            if getattr(shape, "has_text_frame", False):
                yield (
                    location,
                    "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs),
                )
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for cell_index, cell in enumerate(row.cells, start=1):
                        yield (
                            f"{location}/table:{row_index},{cell_index}",
                            "\n".join(paragraph.text for paragraph in cell.text_frame.paragraphs),
                        )
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._text_blocks(shape.shapes, prefix=location)


def format_template_value(value: Any, formatters: tuple[str, ...]) -> str:
    current = value
    for formatter in formatters:
        default_match = DEFAULT_FORMAT.fullmatch(formatter)
        if default_match:
            if current is None or current == "":
                current = default_match.group(1)
            continue
        if current is None or current == "":
            continue
        if formatter == "date":
            if isinstance(current, (date, datetime)):
                current = current.strftime("%Y年%m月%d日")
            else:
                try:
                    current = date.fromisoformat(str(current)[:10]).strftime("%Y年%m月%d日")
                except ValueError:
                    current = str(current)
            continue
        percent_match = PERCENT_FORMAT.fullmatch(formatter)
        if percent_match:
            digits = int(percent_match.group(1) or 2)
            try:
                numeric = Decimal(str(current).strip().rstrip("%"))
                current = f"{numeric:.{digits}f}%"
            except InvalidOperation:
                current = str(current)
    return str("—" if current is None or current == "" else current)
