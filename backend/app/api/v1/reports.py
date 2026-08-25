"""报表模板、自定义定义、合同要素和 PPTX 生成接口。"""

from __future__ import annotations

import io
import os
import uuid
import zipfile
from datetime import date
from pathlib import Path
from typing import Annotated

from fastapi import APIRouter, Depends, File, Form, Request, UploadFile
from fastapi.responses import FileResponse, StreamingResponse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from sqlalchemy import func, select
from sqlalchemy.exc import IntegrityError
from sqlalchemy.orm import Session

from app.api.deps import TenantContext, TenantDatabaseSession, TenantScope, require_roles
from app.api.schemas.reporting import (
    ContractUploadResponse,
    OnlyOfficeSessionResponse,
    ReportBatchCreate,
    ReportBatchItemView,
    ReportBatchView,
    ReportDefinitionCreate,
    ReportDefinitionItem,
    ReportDesignMetadata,
    ReportFieldUpdate,
    ReportFileVersionItem,
    ReportGenerateRequest,
    ReportGenerateResponse,
    ReportLayoutPlacement,
    ReportLayoutUpdate,
    ReportPreviewRequest,
    ReportPreviewResponse,
    ReportProductFieldsResponse,
    ReportRunItem,
    ReportTemplateFromRunRequest,
    ReportTemplateItem,
)
from app.core.config import get_settings
from app.core.credential_security import audit_signing_key
from app.core.errors import AppError
from app.core.files import atomic_write_bytes
from app.db.models import (
    FundProduct,
    ProductDocument,
    ReportBatch,
    ReportBatchItem,
    ReportDefinition,
    ReportFileVersion,
    ReportRun,
    ReportTemplate,
    ReportTemplateVersion,
    UserRole,
)
from app.db.types import utc_now
from app.services.archive_service import sanitize_filename
from app.services.audit_service import AuditService
from app.services.onlyoffice_service import OnlyOfficeService, OnlyOfficeUnavailableError
from app.services.report_field_service import FieldContext, ReportFieldResolver
from app.services.report_presentation_service import ReportPresentationService
from app.services.report_template_service import (
    LEGACY_SNAPSHOT_FIELDS,
    ReportTemplateInspector,
)
from app.services.reporting_service import (
    DEFAULT_REPORT_SECTIONS,
    ReportDataService,
    content_sha256,
    extract_contract_fields,
    extract_contract_text,
)

router = APIRouter()
OperatorScope = Annotated[
    TenantContext,
    Depends(require_roles(UserRole.ADMIN, UserRole.OPERATOR)),
]
_data_service = ReportDataService()
_field_resolver = ReportFieldResolver()
_template_inspector = ReportTemplateInspector()


@router.get("/templates", response_model=list[ReportTemplateItem])
def list_report_templates(
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> list[ReportTemplateItem]:
    del scope
    uploaded = list(
        session.scalars(
            select(ReportTemplate)
            .where(ReportTemplate.is_active.is_(True))
            .order_by(ReportTemplate.name, ReportTemplate.id)
        )
    )
    return [
        ReportTemplateItem(
            key="builtin:weekly",
            name="标准基金周报",
            description="竖版单页周报，包含产品信息、收益指标、净值曲线、策略及合同要素",
            kind="builtin",
            status="builtin",
        ),
        *[
            version_item
            for template in uploaded
            for version_item in _template_items(session, template)
        ],
    ]


@router.post("/templates", response_model=ReportTemplateItem)
async def upload_report_template(
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
    file: Annotated[UploadFile, File(description="带占位符或已知表格结构的 PPTX 模板")],
    name: Annotated[str, Form(min_length=2, max_length=200)],
    description: Annotated[str | None, Form(max_length=1000)] = None,
) -> ReportTemplateItem:
    settings = get_settings()
    filename = file.filename or "report-template.pptx"
    if Path(filename).suffix.casefold() != ".pptx":
        raise AppError("REPORT_TEMPLATE_FORMAT", "报表模板必须是 .pptx 文件")
    content = await file.read(settings.reports.max_template_bytes + 1)
    if len(content) > settings.reports.max_template_bytes:
        raise AppError("REPORT_TEMPLATE_TOO_LARGE", "报表模板超过允许大小")
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise AppError("REPORT_TEMPLATE_INVALID", "PPTX 模板无法打开或结构已损坏") from exc
    inspection = _template_inspector.inspect(presentation, session=session)
    if session.scalar(select(ReportTemplate.id).where(ReportTemplate.name == name.strip())):
        raise AppError("REPORT_TEMPLATE_NAME_EXISTS", "当前租户已存在同名模板", status_code=409)
    safe_name = sanitize_filename(filename)
    relative_path = (
        Path("tenants")
        / str(scope.tenant_id)
        / "reporting"
        / "templates"
        / (f"{uuid.uuid4().hex}_{safe_name}")
    )
    atomic_write_bytes(settings.data_directory / relative_path, content)
    template = ReportTemplate(
        tenant_id=scope.tenant_id,
        name=name.strip(),
        description=description.strip() if description else None,
        original_name=filename,
        stored_path=relative_path.as_posix(),
        content_hash=content_sha256(content),
        is_active=True,
        created_by_user_id=scope.user.id,
    )
    session.add(template)
    session.flush()
    version = ReportTemplateVersion(
        tenant_id=scope.tenant_id,
        template_id=template.id,
        version=1,
        status="draft",
        original_name=filename,
        stored_path=relative_path.as_posix(),
        content_hash=template.content_hash,
        required_fields=list(inspection.required_fields),
        required_components=list(inspection.required_components),
        validation_errors=list(inspection.errors),
        created_by_user_id=scope.user.id,
    )
    session.add(version)
    session.flush()
    _audit(
        session,
        request,
        scope,
        action="report_template.upload",
        resource_type="report_template",
        resource_id=template.id,
        detail={
            "name": template.name,
            "original_name": filename,
            "sha256": template.content_hash,
            "version": 1,
            "required_fields": list(inspection.required_fields),
            "validation_error_count": len(inspection.errors),
        },
    )
    session.commit()
    return _template_item_from_version(template, version)


@router.post("/templates/{template_id}/validate", response_model=ReportTemplateItem)
def validate_report_template(
    template_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportTemplateItem:
    template, version = _editable_template_version(session, template_id)
    version.status = "validating"
    session.flush()
    inspection = _template_inspector.inspect_path(_safe_data_path(version.stored_path), session)
    version.required_fields = list(inspection.required_fields)
    version.required_components = list(inspection.required_components)
    version.validation_errors = list(inspection.errors)
    version.status = "draft"
    _audit(
        session,
        request,
        scope,
        action="report_template.validate",
        resource_type="report_template",
        resource_id=template.id,
        outcome="success" if inspection.is_valid else "failed",
        detail={"version": version.version, "errors": list(inspection.errors)},
    )
    session.commit()
    return _template_item_from_version(template, version)


@router.post("/templates/{template_id}/versions", response_model=ReportTemplateItem)
async def upload_report_template_version(
    template_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
    file: Annotated[UploadFile, File(description="新版 PPTX 模板草稿")],
) -> ReportTemplateItem:
    template = session.get(ReportTemplate, template_id)
    if template is None or not template.is_active:
        raise AppError("REPORT_TEMPLATE_NOT_FOUND", "未找到报表模板", status_code=404)
    existing_draft = session.scalar(
        select(ReportTemplateVersion.id).where(
            ReportTemplateVersion.template_id == template.id,
            ReportTemplateVersion.status.in_(("draft", "validating")),
        )
    )
    if existing_draft is not None:
        raise AppError(
            "REPORT_TEMPLATE_DRAFT_EXISTS",
            "当前模板已有草稿，请先处理现有草稿",
            status_code=409,
        )
    filename = file.filename or "report-template.pptx"
    if Path(filename).suffix.casefold() != ".pptx":
        raise AppError("REPORT_TEMPLATE_FORMAT", "报表模板必须是 .pptx 文件")
    settings = get_settings()
    content = await file.read(settings.reports.max_template_bytes + 1)
    if len(content) > settings.reports.max_template_bytes:
        raise AppError("REPORT_TEMPLATE_TOO_LARGE", "报表模板超过允许大小")
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise AppError("REPORT_TEMPLATE_INVALID", "PPTX 模板无法打开或结构已损坏") from exc
    inspection = _template_inspector.inspect(presentation, session=session)
    next_version = (
        session.scalar(
            select(func.max(ReportTemplateVersion.version)).where(
                ReportTemplateVersion.template_id == template.id
            )
        )
        or 0
    ) + 1
    safe_name = sanitize_filename(filename)
    relative_path = (
        Path("tenants")
        / str(scope.tenant_id)
        / "reporting"
        / "templates"
        / f"{uuid.uuid4().hex}_{safe_name}"
    )
    atomic_write_bytes(settings.data_directory / relative_path, content)
    version = ReportTemplateVersion(
        tenant_id=scope.tenant_id,
        template_id=template.id,
        version=next_version,
        status="draft",
        original_name=filename,
        stored_path=relative_path.as_posix(),
        content_hash=content_sha256(content),
        required_fields=list(inspection.required_fields),
        required_components=list(inspection.required_components),
        validation_errors=list(inspection.errors),
        created_by_user_id=scope.user.id,
    )
    session.add(version)
    session.flush()
    _audit(
        session,
        request,
        scope,
        action="report_template.version.upload",
        resource_type="report_template",
        resource_id=template.id,
        detail={
            "version": next_version,
            "sha256": version.content_hash,
            "validation_error_count": len(inspection.errors),
        },
    )
    session.commit()
    return _template_item_from_version(template, version)


@router.post("/templates/{template_id}/publish", response_model=ReportTemplateItem)
def publish_report_template(
    template_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportTemplateItem:
    template, version = _editable_template_version(session, template_id)
    version.status = "validating"
    session.flush()
    inspection = _template_inspector.inspect_path(_safe_data_path(version.stored_path), session)
    version.required_fields = list(inspection.required_fields)
    version.required_components = list(inspection.required_components)
    version.validation_errors = list(inspection.errors)
    if not inspection.is_valid:
        version.status = "draft"
        session.commit()
        raise AppError(
            "REPORT_TEMPLATE_VALIDATION_FAILED",
            "模板校验失败，请修复占位符后再发布",
            status_code=409,
        )
    for published in session.scalars(
        select(ReportTemplateVersion).where(
            ReportTemplateVersion.template_id == template.id,
            ReportTemplateVersion.status == "published",
        )
    ):
        published.status = "archived"
    version.status = "published"
    version.published_at = utc_now()
    version.published_by_user_id = scope.user.id
    _audit(
        session,
        request,
        scope,
        action="report_template.publish",
        resource_type="report_template",
        resource_id=template.id,
        detail={
            "version": version.version,
            "version_id": version.id,
            "required_fields": version.required_fields,
            "required_components": version.required_components,
        },
    )
    session.commit()
    return _template_item_from_version(template, version)


@router.post(
    "/runs/{run_id}/confirm-template",
    response_model=ReportTemplateItem,
    status_code=201,
)
def confirm_run_as_report_template(
    run_id: int,
    payload: ReportTemplateFromRunRequest,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportTemplateItem:
    """把 OnlyOffice 中已保存的样板报表固化为可批量复用的已发布模板。"""

    run = session.get(ReportRun, run_id)
    if run is None or run.status != "success" or run.current_version_id is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "样板 PPTX 文件不存在", status_code=404)
    version = session.get(ReportFileVersion, run.current_version_id)
    if version is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "样板 PPTX 文件版本不存在", status_code=404)
    name = payload.name.strip()
    if session.scalar(select(ReportTemplate.id).where(ReportTemplate.name == name)):
        raise AppError("REPORT_TEMPLATE_NAME_EXISTS", "当前租户已存在同名模板", status_code=409)

    source_path = _safe_data_path(version.stored_path)
    if not source_path.is_file():
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "样板 PPTX 文件已丢失", status_code=404)
    inspection = _template_inspector.inspect_path(source_path, session)
    if not inspection.tokens:
        raise AppError(
            "REPORT_TEMPLATE_PLACEHOLDER_REQUIRED",
            "模板至少需要一个字段或组件占位符，请在 PPTX 文本框中插入 {{字段键}}",
            status_code=409,
        )
    if not inspection.is_valid:
        raise AppError(
            "REPORT_TEMPLATE_VALIDATION_FAILED",
            "模板占位符校验失败，请修复后再确认",
            status_code=409,
        )

    settings = get_settings()
    content = source_path.read_bytes()
    safe_name = sanitize_filename(f"{name}.pptx")
    relative_path = (
        Path("tenants")
        / str(scope.tenant_id)
        / "reporting"
        / "templates"
        / f"{uuid.uuid4().hex}_{safe_name}"
    )
    atomic_write_bytes(settings.data_directory / relative_path, content)
    digest = content_sha256(content)
    template = ReportTemplate(
        tenant_id=scope.tenant_id,
        name=name,
        description=payload.description.strip() if payload.description else None,
        original_name=safe_name,
        stored_path=relative_path.as_posix(),
        content_hash=digest,
        is_active=True,
        created_by_user_id=scope.user.id,
    )
    session.add(template)
    session.flush()
    template_version = ReportTemplateVersion(
        tenant_id=scope.tenant_id,
        template_id=template.id,
        version=1,
        status="published",
        original_name=safe_name,
        stored_path=relative_path.as_posix(),
        content_hash=digest,
        required_fields=list(inspection.required_fields),
        required_components=list(inspection.required_components),
        validation_errors=[],
        published_at=utc_now(),
        created_by_user_id=scope.user.id,
        published_by_user_id=scope.user.id,
    )
    session.add(template_version)
    session.flush()
    _audit(
        session,
        request,
        scope,
        action="report_template.confirm_from_run",
        resource_type="report_template",
        resource_id=template.id,
        detail={
            "source_run_id": run.id,
            "source_file_version_id": version.id,
            "version_id": template_version.id,
            "required_fields": template_version.required_fields,
            "required_components": template_version.required_components,
            "sha256": digest,
        },
    )
    session.commit()
    return _template_item_from_version(template, template_version)


@router.get("/runs/{run_id}/design-metadata", response_model=ReportDesignMetadata)
def get_report_design_metadata(
    run_id: int,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportDesignMetadata:
    del scope
    run, version, path = _editable_report_output(session, run_id)
    presentation = Presentation(str(path))
    placements: list[ReportLayoutPlacement] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if not shape.name.startswith("codex-field:") or not getattr(
                shape, "has_text_frame", False
            ):
                continue
            placement_id = shape.name.removeprefix("codex-field:")
            placements.append(
                ReportLayoutPlacement(
                    id=placement_id,
                    token=shape.text.strip(),
                    slide=slide_number,
                    x=shape.left / presentation.slide_width,
                    y=shape.top / presentation.slide_height,
                    width=shape.width / presentation.slide_width,
                    height=shape.height / presentation.slide_height,
                )
            )
    return ReportDesignMetadata(
        slide_count=len(presentation.slides),
        slide_width=presentation.slide_width,
        slide_height=presentation.slide_height,
        placements=placements,
    )


@router.post("/runs/{run_id}/layout", response_model=ReportRunItem)
def update_report_layout(
    run_id: int,
    payload: ReportLayoutUpdate,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportRunItem:
    """按归一化坐标把字段占位符写入 PPTX，并追加不可变文件版本。"""

    run, current, path = _editable_report_output(session, run_id)
    presentation = Presentation(str(path))
    for slide in presentation.slides:
        for shape in list(slide.shapes):
            if shape.name.startswith("codex-field:"):
                shape._element.getparent().remove(shape._element)  # noqa: SLF001
    for item in payload.placements:
        if item.slide > len(presentation.slides):
            raise AppError("REPORT_LAYOUT_SLIDE_INVALID", "字段绑定的幻灯片页码不存在")
        slide = presentation.slides[item.slide - 1]
        textbox = slide.shapes.add_textbox(
            int(item.x * presentation.slide_width),
            int(item.y * presentation.slide_height),
            int(item.width * presentation.slide_width),
            int(item.height * presentation.slide_height),
        )
        textbox.name = f"codex-field:{item.id}"
        textbox.text_frame.clear()
        paragraph = textbox.text_frame.paragraphs[0]
        run_text = paragraph.add_run()
        run_text.text = item.token
        run_text.font.size = Pt(item.font_size)
        run_text.font.bold = item.bold
        run_text.font.color.rgb = RGBColor.from_string(item.color.removeprefix("#"))

    output = io.BytesIO()
    presentation.save(output)
    content = output.getvalue()
    next_version = (
        session.scalar(
            select(func.max(ReportFileVersion.version)).where(
                ReportFileVersion.report_run_id == run.id
            )
        )
        or 0
    ) + 1
    filename = run.output_filename or current.filename
    relative_path = (
        Path("tenants")
        / str(scope.tenant_id)
        / "reporting"
        / "exports"
        / f"{run.report_date.year:04d}"
        / f"{run.report_date.month:02d}"
        / str(run.id)
        / f"v{next_version}_{uuid.uuid4().hex[:8]}_{sanitize_filename(filename)}"
    )
    atomic_write_bytes(get_settings().data_directory / relative_path, content)
    version = ReportFileVersion(
        tenant_id=scope.tenant_id,
        report_run_id=run.id,
        version=next_version,
        source="designer",
        filename=filename,
        stored_path=relative_path.as_posix(),
        content_hash=content_sha256(content),
        file_size=len(content),
        created_by_user_id=scope.user.id,
    )
    session.add(version)
    session.flush()
    run.current_version_id = version.id
    run.output_path = version.stored_path
    _audit(
        session,
        request,
        scope,
        action="report.layout.update",
        resource_type="report_file_version",
        resource_id=version.id,
        detail={"run_id": run.id, "version": next_version, "placements": len(payload.placements)},
    )
    session.commit()
    product = session.get(FundProduct, run.fund_product_id)
    return _run_item(run, product.product_name if product else str(run.fund_product_id), session)


@router.get("/product-fields/{product_id}", response_model=ReportProductFieldsResponse)
def get_report_product_fields(
    product_id: int,
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> ReportProductFieldsResponse:
    del scope
    product = _data_service.get_product(session, product_id)
    return ReportProductFieldsResponse(
        product_id=product.id,
        product_code=product.product_code,
        product_name=product.product_name,
        fields=_data_service.fields(session, product),
    )


@router.patch(
    "/product-fields/{product_id}/{field_key}", response_model=ReportProductFieldsResponse
)
def update_report_product_field(
    product_id: int,
    field_key: str,
    payload: ReportFieldUpdate,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportProductFieldsResponse:
    product = _data_service.get_product(session, product_id)
    before, after = _data_service.update_field(
        product,
        field_key=field_key,
        value=payload.value,
        restore_source=payload.restore_source,
    )
    _audit(
        session,
        request,
        scope,
        action=(
            "report_product_field.restore"
            if payload.restore_source
            else "report_product_field.update"
        ),
        resource_type="fund_product",
        resource_id=product.id,
        detail={
            "field_key": field_key,
            "before": before,
            "after": after,
            "reason": payload.reason,
        },
    )
    session.commit()
    return ReportProductFieldsResponse(
        product_id=product.id,
        product_code=product.product_code,
        product_name=product.product_name,
        fields=_data_service.fields(session, product),
    )


@router.post("/contracts/{product_id}", response_model=ContractUploadResponse)
async def upload_product_contract(
    product_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
    file: Annotated[UploadFile, File(description="PDF、DOCX 或 TXT 产品合同")],
) -> ContractUploadResponse:
    settings = get_settings()
    product = _data_service.get_product(session, product_id)
    filename = file.filename or "contract.pdf"
    content = await file.read(settings.reports.max_contract_bytes + 1)
    if len(content) > settings.reports.max_contract_bytes:
        raise AppError("CONTRACT_TOO_LARGE", "合同文件超过允许大小")
    digest = content_sha256(content)
    existing = session.scalar(
        select(ProductDocument).where(
            ProductDocument.fund_product_id == product.id,
            ProductDocument.content_hash == digest,
        )
    )
    if existing is not None:
        raise AppError("CONTRACT_ALREADY_ARCHIVED", "该合同已经归档", status_code=409)
    text = extract_contract_text(filename, content)
    fields = extract_contract_fields(text)
    safe_name = sanitize_filename(filename)
    today = date.today()
    relative_path = (
        Path("tenants")
        / str(scope.tenant_id)
        / "reporting"
        / "contracts"
        / product.product_code
        / f"{today.year:04d}"
        / f"{today.month:02d}"
        / f"{uuid.uuid4().hex}_{safe_name}"
    )
    atomic_write_bytes(settings.data_directory / relative_path, content)
    document = ProductDocument(
        tenant_id=scope.tenant_id,
        fund_product_id=product.id,
        document_type="contract",
        original_name=filename,
        stored_path=relative_path.as_posix(),
        content_hash=digest,
        extracted_fields=fields,
        created_by_user_id=scope.user.id,
    )
    session.add(document)
    session.flush()
    _data_service.apply_contract_fields(product, fields, document_id=document.id, filename=filename)
    _audit(
        session,
        request,
        scope,
        action="product_contract.upload",
        resource_type="product_document",
        resource_id=document.id,
        detail={
            "product_id": product.id,
            "original_name": filename,
            "sha256": digest,
            "extracted_fields": sorted(fields),
        },
    )
    session.commit()
    return ContractUploadResponse(
        document_id=document.id,
        original_name=filename,
        extracted_fields=fields,
        extracted_count=len(fields),
    )


@router.get("/definitions", response_model=list[ReportDefinitionItem])
def list_report_definitions(
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> list[ReportDefinitionItem]:
    del scope
    return [
        ReportDefinitionItem.model_validate(item, from_attributes=True)
        for item in session.scalars(
            select(ReportDefinition).order_by(ReportDefinition.update_time.desc())
        )
    ]


@router.post("/definitions", response_model=ReportDefinitionItem)
def create_report_definition(
    payload: ReportDefinitionCreate,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportDefinitionItem:
    _data_service.get_product(session, payload.fund_product_id)
    _resolve_template(session, payload.template_key)
    existing_id = session.scalar(
        select(ReportDefinition.id).where(ReportDefinition.name == payload.name.strip())
    )
    if existing_id:
        raise AppError("REPORT_DEFINITION_NAME_EXISTS", "当前租户已存在同名报表", status_code=409)
    sections = _validated_sections(payload.sections)
    definition = ReportDefinition(
        tenant_id=scope.tenant_id,
        name=payload.name.strip(),
        fund_product_id=payload.fund_product_id,
        template_key=payload.template_key,
        report_type=payload.report_type,
        sections=sections,
        settings=payload.settings,
        created_by_user_id=scope.user.id,
    )
    session.add(definition)
    session.flush()
    _audit(
        session,
        request,
        scope,
        action="report_definition.create",
        resource_type="report_definition",
        resource_id=definition.id,
        detail={
            "name": definition.name,
            "template_key": definition.template_key,
            "sections": sections,
        },
    )
    session.commit()
    return ReportDefinitionItem.model_validate(definition, from_attributes=True)


@router.post("/preview", response_model=ReportPreviewResponse)
def preview_report(
    payload: ReportPreviewRequest,
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> ReportPreviewResponse:
    del scope
    product = _data_service.get_product(session, payload.fund_product_id)
    snapshot = _data_service.build_snapshot(
        session,
        product,
        report_date=payload.report_date,
        share_product_code=_share_code(payload.settings),
    )
    return ReportPreviewResponse.model_validate(snapshot)


@router.post("/generate", response_model=ReportGenerateResponse)
def generate_report(
    payload: ReportGenerateRequest,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportGenerateResponse:
    definition = None
    settings_payload = dict(payload.settings)
    if payload.definition_id is not None:
        definition = session.get(ReportDefinition, payload.definition_id)
        if definition is None:
            raise AppError("REPORT_DEFINITION_NOT_FOUND", "未找到报表定义", status_code=404)
        product_id = definition.fund_product_id
        template_key = definition.template_key
        sections = list(definition.sections)
        settings_payload = {**(definition.settings or {}), **settings_payload}
    else:
        product_id = payload.fund_product_id
        template_key = payload.template_key or "builtin:weekly"
        sections = payload.sections or list(DEFAULT_REPORT_SECTIONS)
    assert product_id is not None
    product = _data_service.get_product(session, product_id)
    template = _resolve_template(session, template_key)
    sections = _validated_sections(sections)
    snapshot = _data_service.build_snapshot(
        session,
        product,
        report_date=payload.report_date,
        share_product_code=_share_code(settings_payload),
    )
    if isinstance(template, ReportTemplateVersion):
        dynamic_keys = [
            key for key in (template.required_fields or []) if key not in LEGACY_SNAPSHOT_FIELDS
        ]
        if dynamic_keys:
            resolved = _field_resolver.resolve_many(
                session,
                dynamic_keys,
                FieldContext(
                    tenant_id=scope.tenant_id,
                    tenant_name=scope.tenant_name,
                    product_id=product.id,
                    report_date=date.fromisoformat(snapshot["report_date"]),
                ),
                allow_inactive=True,
            )
            snapshot["dynamic_fields"] = {key: value.value for key, (_, value) in resolved.items()}
            snapshot["dynamic_field_provenance"] = {
                key: {
                    "source_type": value.source_type,
                    "source_reference": value.source_reference,
                    "used_default": value.used_default,
                }
                for key, (_, value) in resolved.items()
            }
            snapshot["field_definition_versions"] = {
                key: definition["version"] for key, (definition, _) in resolved.items()
            }
    template_version_id = template.id if isinstance(template, ReportTemplateVersion) else None
    snapshot["template_version_id"] = template_version_id
    snapshot["template_content_hash"] = (
        template.content_hash if isinstance(template, ReportTemplateVersion) else None
    )
    actual_date = date.fromisoformat(snapshot["report_date"])
    run = ReportRun(
        tenant_id=scope.tenant_id,
        definition_id=definition.id if definition else None,
        fund_product_id=product.id,
        template_key=template_key,
        template_version_id=template_version_id,
        report_date=actual_date,
        status="processing",
        input_snapshot={**snapshot, "sections": sections, "settings": settings_payload},
        field_definition_versions=dict(snapshot.get("field_definition_versions") or {}),
        created_by_user_id=scope.user.id,
    )
    session.add(run)
    session.flush()
    safe_product = sanitize_filename(product.product_name, max_length=100)
    filename = f"{safe_product}_{actual_date.isoformat()}_基金周报.pptx"
    template_path = _stored_template_path(template) if template else None
    try:
        _create_file_version(
            session,
            run,
            snapshot,
            sections,
            template_path,
            filename,
            source="generated",
            user_id=scope.user.id,
        )
        run.status = "success"
        outcome = "success"
    except ReportGenerationStageError as exc:
        run.status = "failed"
        run.error_stage = exc.stage
        run.error_code = exc.code
        run.error_message = str(exc)[:1000]
        outcome = "failed"
    _audit(
        session,
        request,
        scope,
        action="report.generate",
        resource_type="report_run",
        resource_id=run.id,
        outcome=outcome,
        detail={
            "product_id": product.id,
            "template_key": template_key,
            "report_date": actual_date.isoformat(),
            "sections": sections,
            "output_filename": run.output_filename,
            "error": run.error_message,
            "error_stage": run.error_stage,
            "error_code": run.error_code,
        },
    )
    session.commit()
    item = _run_item(run, product.product_name, session)
    if run.status != "success":
        raise AppError("REPORT_GENERATION_FAILED", "报表生成失败，请查看生成记录", status_code=500)
    return ReportGenerateResponse(run=item, download_url=f"/api/v1/reports/runs/{run.id}/download")


@router.get("/runs", response_model=list[ReportRunItem])
def list_report_runs(
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> list[ReportRunItem]:
    del scope
    rows = session.execute(
        select(ReportRun, FundProduct.product_name)
        .join(FundProduct, FundProduct.id == ReportRun.fund_product_id)
        .order_by(ReportRun.id.desc())
        .limit(100)
    )
    return [_run_item(run, product_name, session) for run, product_name in rows]


@router.get("/runs/{run_id}/download")
def download_report(
    run_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> FileResponse:
    run = session.get(ReportRun, run_id)
    if run is None or run.status != "success" or run.current_version_id is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件不存在", status_code=404)
    version = session.get(ReportFileVersion, run.current_version_id)
    if version is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件版本不存在", status_code=404)
    path = _safe_data_path(version.stored_path)
    if not path.is_file():
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件已丢失", status_code=404)
    _audit(
        session,
        request,
        scope,
        action="report.download",
        resource_type="report_file_version",
        resource_id=version.id,
        detail={"run_id": run.id, "version": version.version, "sha256": version.content_hash},
    )
    session.commit()
    return FileResponse(
        path,
        filename=version.filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@router.get("/runs/{run_id}/versions", response_model=list[ReportFileVersionItem])
def list_report_file_versions(
    run_id: int,
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> list[ReportFileVersionItem]:
    del scope
    if session.get(ReportRun, run_id) is None:
        raise AppError("REPORT_RUN_NOT_FOUND", "报表生成记录不存在", status_code=404)
    return [
        ReportFileVersionItem.model_validate(version, from_attributes=True)
        for version in session.scalars(
            select(ReportFileVersion)
            .where(ReportFileVersion.report_run_id == run_id)
            .order_by(ReportFileVersion.version.desc())
        )
    ]


@router.post(
    "/runs/{run_id}/onlyoffice/session",
    response_model=OnlyOfficeSessionResponse,
)
def create_onlyoffice_view_session(
    run_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> OnlyOfficeSessionResponse:
    run = session.get(ReportRun, run_id)
    if run is None or run.status != "success" or run.current_version_id is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件不存在", status_code=404)
    version = session.get(ReportFileVersion, run.current_version_id)
    if version is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件版本不存在", status_code=404)
    path = _safe_data_path(version.stored_path)
    if not path.is_file():
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件已丢失", status_code=404)
    service = OnlyOfficeService(get_settings().onlyoffice)
    try:
        service.ensure_ready()
    except OnlyOfficeUnavailableError as exc:
        raise AppError(
            "ONLYOFFICE_UNAVAILABLE",
            str(exc),
            status_code=503,
        ) from exc
    editable = scope.role in (UserRole.ADMIN, UserRole.OPERATOR)
    payload = service.build_session(
        tenant_id=scope.tenant_id,
        run_id=run.id,
        version_id=version.id,
        content_hash=version.content_hash,
        filename=version.filename,
        user_id=scope.user.id,
        username=scope.user.username,
        editable=editable,
    )
    _audit(
        session,
        request,
        scope,
        action="report.onlyoffice.edit" if editable else "report.onlyoffice.view",
        resource_type="report_file_version",
        resource_id=version.id,
        detail={
            "run_id": run.id,
            "version": version.version,
            "mode": "edit" if editable else "view",
        },
    )
    session.commit()
    return OnlyOfficeSessionResponse.model_validate(payload)


@router.get("/runs/{run_id}/versions/{version_id}/download")
def download_report_file_version(
    run_id: int,
    version_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> FileResponse:
    version = session.get(ReportFileVersion, version_id)
    if version is None or version.report_run_id != run_id:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件版本不存在", status_code=404)
    path = _safe_data_path(version.stored_path)
    if not path.is_file():
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "报表文件已丢失", status_code=404)
    _audit(
        session,
        request,
        scope,
        action="report.version.download",
        resource_type="report_file_version",
        resource_id=version.id,
        detail={"run_id": run_id, "version": version.version, "sha256": version.content_hash},
    )
    session.commit()
    return FileResponse(
        path,
        filename=version.filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@router.post("/runs/{run_id}/regenerate", response_model=ReportGenerateResponse)
def regenerate_report_from_snapshot(
    run_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportGenerateResponse:
    run = session.get(ReportRun, run_id)
    if run is None:
        raise AppError("REPORT_RUN_NOT_FOUND", "报表生成记录不存在", status_code=404)
    product = session.get(FundProduct, run.fund_product_id)
    if product is None:
        raise AppError("FUND_PRODUCT_NOT_FOUND", "原基金产品不存在", status_code=404)
    snapshot = dict(run.input_snapshot or {})
    sections = _validated_sections(list(snapshot.get("sections") or []))
    template = (
        session.get(ReportTemplateVersion, run.template_version_id)
        if run.template_version_id
        else None
    )
    if run.template_version_id and template is None:
        raise AppError("REPORT_TEMPLATE_FILE_MISSING", "原模板版本不存在", status_code=404)
    template_path = _stored_template_path(template) if template else None
    filename = run.output_filename or (
        f"{sanitize_filename(product.product_name, max_length=100)}_"
        f"{run.report_date.isoformat()}_基金周报.pptx"
    )
    run.status = "processing"
    run.error_stage = run.error_code = run.error_message = None
    try:
        _create_file_version(
            session,
            run,
            snapshot,
            sections,
            template_path,
            filename,
            source="regenerated",
            user_id=scope.user.id,
        )
        run.status = "success"
        outcome = "success"
    except ReportGenerationStageError as exc:
        run.status = "failed"
        run.error_stage = exc.stage
        run.error_code = exc.code
        run.error_message = str(exc)[:1000]
        outcome = "failed"
    _audit(
        session,
        request,
        scope,
        action="report.regenerate",
        resource_type="report_run",
        resource_id=run.id,
        outcome=outcome,
        detail={
            "current_version_id": run.current_version_id,
            "error_stage": run.error_stage,
            "error_code": run.error_code,
        },
    )
    session.commit()
    if run.status != "success":
        raise AppError("REPORT_GENERATION_FAILED", "报表重新生成失败", status_code=500)
    return ReportGenerateResponse(
        run=_run_item(run, product.product_name, session),
        download_url=f"/api/v1/reports/runs/{run.id}/download",
    )


@router.post("/batches", response_model=ReportBatchView, status_code=201)
def create_report_batch(
    payload: ReportBatchCreate,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportBatchView:
    existing = session.scalar(
        select(ReportBatch).where(ReportBatch.idempotency_key == payload.idempotency_key)
    )
    if existing is not None:
        return _batch_view(existing)
    product_ids = list(dict.fromkeys(payload.product_ids))
    product_query = select(FundProduct)
    if product_ids:
        product_query = product_query.where(FundProduct.id.in_(product_ids))
    if payload.product_code_contains:
        product_query = product_query.where(
            FundProduct.product_code.ilike(f"%{payload.product_code_contains.strip()}%")
        )
    if payload.product_name_contains:
        product_query = product_query.where(
            FundProduct.product_name.ilike(f"%{payload.product_name_contains.strip()}%")
        )
    products = list(session.scalars(product_query.order_by(FundProduct.id).limit(1001)))
    if product_ids and len(products) != len(product_ids):
        raise AppError(
            "REPORT_BATCH_PRODUCT_NOT_FOUND", "部分基金不存在或无权访问", status_code=404
        )
    if not products:
        raise AppError("REPORT_BATCH_PRODUCT_EMPTY", "筛选条件没有匹配基金")
    if len(products) > 1000:
        raise AppError("REPORT_BATCH_TOO_LARGE", "单个批次最多包含 1000 只基金")
    template = _resolve_template(session, payload.template_key)
    if template is not None and not isinstance(template, ReportTemplateVersion):
        raise AppError(
            "REPORT_TEMPLATE_NOT_PUBLISHED",
            "批量生成只能使用已发布模板版本",
            status_code=409,
        )
    template_version_id = template.id if isinstance(template, ReportTemplateVersion) else None
    sections = _validated_sections(payload.sections)
    batch = ReportBatch(
        tenant_id=scope.tenant_id,
        idempotency_key=payload.idempotency_key,
        template_version_id=template_version_id,
        template_key=payload.template_key,
        report_date=payload.report_date,
        sections=sections,
        settings=payload.settings,
        status="pending",
        total_count=len(products),
        created_by_user_id=scope.user.id,
    )
    session.add(batch)
    session.flush()
    for product in products:
        snapshot = _data_service.build_snapshot(
            session,
            product,
            report_date=payload.report_date,
            share_product_code=_share_code(payload.settings),
        )
        if isinstance(template, ReportTemplateVersion):
            keys = [
                key for key in (template.required_fields or []) if key not in LEGACY_SNAPSHOT_FIELDS
            ]
            if keys:
                resolved = _field_resolver.resolve_many(
                    session,
                    keys,
                    FieldContext(
                        scope.tenant_id,
                        scope.tenant_name,
                        product.id,
                        payload.report_date,
                    ),
                    allow_inactive=True,
                )
                snapshot["dynamic_fields"] = {
                    key: value.value for key, (_, value) in resolved.items()
                }
                snapshot["dynamic_field_provenance"] = {
                    key: {
                        "source_type": value.source_type,
                        "source_reference": value.source_reference,
                        "used_default": value.used_default,
                    }
                    for key, (_, value) in resolved.items()
                }
                snapshot["field_definition_versions"] = {
                    key: definition["version"] for key, (definition, _) in resolved.items()
                }
        snapshot.update(
            {
                "sections": sections,
                "settings": payload.settings,
                "template_version_id": template_version_id,
                "template_content_hash": getattr(template, "content_hash", None),
            }
        )
        session.add(
            ReportBatchItem(
                tenant_id=scope.tenant_id,
                batch_id=batch.id,
                fund_product_id=product.id,
                status="pending",
                idempotency_key=f"{scope.tenant_id}:{payload.idempotency_key}:{product.id}",
                input_snapshot=snapshot,
            )
        )
    _audit(
        session,
        request,
        scope,
        action="report_batch.create",
        resource_type="report_batch",
        resource_id=batch.id,
        detail={"total_count": len(products), "template_version_id": template_version_id},
    )
    try:
        session.commit()
    except IntegrityError:
        session.rollback()
        concurrent = session.scalar(
            select(ReportBatch).where(ReportBatch.idempotency_key == payload.idempotency_key)
        )
        if concurrent is None:
            raise
        return _batch_view(concurrent)
    return _batch_view(batch)


@router.get("/batches/{batch_id}", response_model=ReportBatchView)
def get_report_batch(
    batch_id: int, session: TenantDatabaseSession, scope: TenantScope
) -> ReportBatchView:
    del scope
    batch = session.get(ReportBatch, batch_id)
    if batch is None:
        raise AppError("REPORT_BATCH_NOT_FOUND", "批次不存在", status_code=404)
    return _batch_view(batch)


@router.get("/batches/{batch_id}/items", response_model=list[ReportBatchItemView])
def get_report_batch_items(
    batch_id: int, session: TenantDatabaseSession, scope: TenantScope
) -> list[ReportBatchItemView]:
    del scope
    if session.get(ReportBatch, batch_id) is None:
        raise AppError("REPORT_BATCH_NOT_FOUND", "批次不存在", status_code=404)
    rows = session.execute(
        select(ReportBatchItem, FundProduct.product_name)
        .join(FundProduct, FundProduct.id == ReportBatchItem.fund_product_id)
        .where(ReportBatchItem.batch_id == batch_id)
        .order_by(ReportBatchItem.id)
    )
    return [_batch_item_view(item, name) for item, name in rows]


@router.post("/batches/{batch_id}/retry", response_model=ReportBatchView)
def retry_report_batch(
    batch_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportBatchView:
    batch = session.get(ReportBatch, batch_id)
    if batch is None:
        raise AppError("REPORT_BATCH_NOT_FOUND", "批次不存在", status_code=404)
    items = list(
        session.scalars(
            select(ReportBatchItem).where(
                ReportBatchItem.batch_id == batch_id,
                ReportBatchItem.status == "failed",
            )
        )
    )
    for item in items:
        item.status = "pending"
        item.error_code = item.error_message = item.locked_by = item.locked_at = None
    batch.status = "pending"
    batch.failed_count = max(0, batch.failed_count - len(items))
    _audit(
        session,
        request,
        scope,
        action="report_batch.retry",
        resource_type="report_batch",
        resource_id=batch.id,
        detail={"retry_count": len(items)},
    )
    session.commit()
    return _batch_view(batch)


@router.post("/batches/{batch_id}/cancel", response_model=ReportBatchView)
def cancel_report_batch(
    batch_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: OperatorScope,
) -> ReportBatchView:
    batch = session.get(ReportBatch, batch_id)
    if batch is None:
        raise AppError("REPORT_BATCH_NOT_FOUND", "批次不存在", status_code=404)
    pending = list(
        session.scalars(
            select(ReportBatchItem).where(
                ReportBatchItem.batch_id == batch_id,
                ReportBatchItem.status == "pending",
            )
        )
    )
    for item in pending:
        item.status = "cancelled"
    batch.cancelled_count += len(pending)
    batch.status = "cancelled"
    _audit(
        session,
        request,
        scope,
        action="report_batch.cancel",
        resource_type="report_batch",
        resource_id=batch.id,
        detail={"cancelled_count": len(pending)},
    )
    session.commit()
    return _batch_view(batch)


@router.get("/batches/{batch_id}/download")
def download_report_batch(
    batch_id: int,
    request: Request,
    session: TenantDatabaseSession,
    scope: TenantScope,
) -> StreamingResponse:
    batch = session.get(ReportBatch, batch_id)
    if batch is None:
        raise AppError("REPORT_BATCH_NOT_FOUND", "批次不存在", status_code=404)
    items = list(
        session.scalars(
            select(ReportBatchItem)
            .where(ReportBatchItem.batch_id == batch_id)
            .order_by(ReportBatchItem.id)
        )
    )
    content = io.BytesIO()
    failures = ["product_id,error_code,error_message"]
    with zipfile.ZipFile(content, "w", zipfile.ZIP_DEFLATED) as archive:
        for item in items:
            if item.status == "success" and item.report_run_id:
                run = session.get(ReportRun, item.report_run_id)
                version = (
                    session.get(ReportFileVersion, run.current_version_id)
                    if run and run.current_version_id
                    else None
                )
                if version:
                    path = _safe_data_path(version.stored_path)
                    if path.is_file():
                        archive.write(path, arcname=f"{item.id}_{version.filename}")
                        continue
            if item.status in {"failed", "cancelled"}:
                message = (item.error_message or "").replace('"', "'")
                failures.append(
                    f'{item.fund_product_id},{item.error_code or item.status},"{message}"'
                )
        archive.writestr("失败清单.csv", "\ufeff" + "\n".join(failures))
    _audit(
        session,
        request,
        scope,
        action="report_batch.download",
        resource_type="report_batch",
        resource_id=batch.id,
        detail={"success_count": batch.success_count, "failed_count": batch.failed_count},
    )
    session.commit()
    content.seek(0)
    return StreamingResponse(
        content,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="report-batch-{batch.id}.zip"'},
    )


class ReportGenerationStageError(RuntimeError):
    def __init__(self, stage: str, code: str, message: str) -> None:
        super().__init__(message)
        self.stage = stage
        self.code = code


def _create_file_version(
    session: Session,
    run: ReportRun,
    snapshot: dict,
    sections: list[str],
    template_path: Path | None,
    filename: str,
    *,
    source: str,
    user_id: int | None,
) -> ReportFileVersion:
    next_version = (
        session.scalar(
            select(func.max(ReportFileVersion.version)).where(
                ReportFileVersion.report_run_id == run.id
            )
        )
        or 0
    ) + 1
    relative_path = (
        Path("tenants")
        / str(run.tenant_id)
        / "reporting"
        / "exports"
        / f"{run.report_date.year:04d}"
        / f"{run.report_date.month:02d}"
        / str(run.id)
        / f"v{next_version}_{filename}"
    )
    final_path = get_settings().data_directory / relative_path
    final_path.parent.mkdir(parents=True, exist_ok=True)
    temp_path = final_path.with_name(f".{final_path.name}.{uuid.uuid4().hex}.tmp.pptx")
    try:
        try:
            ReportPresentationService().generate(
                snapshot,
                output_path=temp_path,
                sections=sections,
                template_path=template_path,
            )
        except Exception as exc:
            raise ReportGenerationStageError("render", "REPORT_RENDER_FAILED", str(exc)) from exc
        try:
            rendered = Presentation(temp_path)
            if not rendered.slides:
                raise ValueError("生成文件没有幻灯片")
        except Exception as exc:
            raise ReportGenerationStageError(
                "validate", "REPORT_OUTPUT_INVALID", "生成的 PPTX 无法重新打开"
            ) from exc
        try:
            content = temp_path.read_bytes()
            digest = content_sha256(content)
            os.replace(temp_path, final_path)
        except Exception as exc:
            raise ReportGenerationStageError(
                "store", "REPORT_OUTPUT_STORE_FAILED", "报表文件保存失败"
            ) from exc
        version = ReportFileVersion(
            tenant_id=run.tenant_id,
            report_run_id=run.id,
            version=next_version,
            source=source,
            filename=filename,
            stored_path=relative_path.as_posix(),
            content_hash=digest,
            file_size=len(content),
            created_by_user_id=user_id,
        )
        session.add(version)
        session.flush()
        run.current_version_id = version.id
        run.output_filename = filename
        run.output_path = relative_path.as_posix()
        return version
    finally:
        temp_path.unlink(missing_ok=True)


def _resolve_template(
    session: Session, template_key: str
) -> ReportTemplate | ReportTemplateVersion | None:
    if template_key == "builtin:weekly":
        return None
    if template_key.startswith("template-version:"):
        try:
            version_id = int(template_key.split(":", 1)[1])
        except ValueError as exc:
            raise AppError(
                "REPORT_TEMPLATE_NOT_FOUND", "模板版本编号无效", status_code=404
            ) from exc
        version = session.get(ReportTemplateVersion, version_id)
        if version is None or version.status != "published":
            raise AppError("REPORT_TEMPLATE_NOT_PUBLISHED", "模板版本未发布", status_code=409)
        return version
    if not template_key.startswith("uploaded:"):
        raise AppError("REPORT_TEMPLATE_NOT_FOUND", "未找到报表模板", status_code=404)
    try:
        template_id = int(template_key.split(":", 1)[1])
    except ValueError as exc:
        raise AppError("REPORT_TEMPLATE_NOT_FOUND", "报表模板编号无效", status_code=404) from exc
    template = session.get(ReportTemplate, template_id)
    if template is None or not template.is_active:
        raise AppError("REPORT_TEMPLATE_NOT_FOUND", "未找到报表模板", status_code=404)
    published = session.scalar(
        select(ReportTemplateVersion)
        .where(
            ReportTemplateVersion.template_id == template.id,
            ReportTemplateVersion.status == "published",
        )
        .order_by(ReportTemplateVersion.version.desc())
    )
    return published or template


def _stored_template_path(template: ReportTemplate | ReportTemplateVersion) -> Path:
    path = _safe_data_path(template.stored_path)
    if not path.is_file():
        raise AppError("REPORT_TEMPLATE_FILE_MISSING", "模板文件已丢失")
    return path


def _template_items(session: Session, template: ReportTemplate) -> list[ReportTemplateItem]:
    versions = list(
        session.scalars(
            select(ReportTemplateVersion)
            .where(
                ReportTemplateVersion.template_id == template.id,
                ReportTemplateVersion.status.in_(("draft", "validating", "published")),
            )
            .order_by(ReportTemplateVersion.version.desc())
        )
    )
    latest_by_status: dict[str, ReportTemplateVersion] = {}
    for version in versions:
        latest_by_status.setdefault(version.status, version)
    if not latest_by_status:
        return [
            ReportTemplateItem(
                key=f"uploaded:{template.id}",
                id=template.id,
                name=template.name,
                description=template.description,
                kind="uploaded",
                original_name=template.original_name,
                is_active=template.is_active,
                create_time=template.create_time,
            )
        ]
    return [
        _template_item_from_version(template, version)
        for status in ("published", "validating", "draft")
        if (version := latest_by_status.get(status)) is not None
    ]


def _template_item_from_version(
    template: ReportTemplate, version: ReportTemplateVersion
) -> ReportTemplateItem:
    key = (
        f"template-version:{version.id}" if version.status == "published" else f"draft:{version.id}"
    )
    return ReportTemplateItem(
        key=key,
        id=template.id,
        name=template.name,
        description=template.description,
        kind="uploaded",
        original_name=version.original_name,
        is_active=template.is_active,
        create_time=version.create_time,
        version_id=version.id,
        version=version.version,
        status=version.status,
        required_fields=list(version.required_fields or []),
        required_components=list(version.required_components or []),
        validation_errors=list(version.validation_errors or []),
    )


def _editable_template_version(
    session: Session, template_id: int
) -> tuple[ReportTemplate, ReportTemplateVersion]:
    template = session.get(ReportTemplate, template_id)
    if template is None or not template.is_active:
        raise AppError("REPORT_TEMPLATE_NOT_FOUND", "未找到报表模板", status_code=404)
    version = session.scalar(
        select(ReportTemplateVersion)
        .where(
            ReportTemplateVersion.template_id == template.id,
            ReportTemplateVersion.status.in_(("draft", "validating")),
        )
        .order_by(ReportTemplateVersion.version.desc())
    )
    if version is None:
        raise AppError("REPORT_TEMPLATE_DRAFT_NOT_FOUND", "当前模板没有可编辑草稿", status_code=409)
    return template, version


def _safe_data_path(relative_path: str) -> Path:
    root = get_settings().data_directory.resolve()
    path = (root / relative_path).resolve()
    if path != root and root not in path.parents:
        raise AppError("REPORT_PATH_INVALID", "报表文件路径无效")
    return path


def _editable_report_output(
    session: Session, run_id: int
) -> tuple[ReportRun, ReportFileVersion, Path]:
    run = session.get(ReportRun, run_id)
    if run is None or run.status != "success" or run.current_version_id is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "PPTX 样板不存在", status_code=404)
    version = session.get(ReportFileVersion, run.current_version_id)
    if version is None:
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "PPTX 样板版本不存在", status_code=404)
    path = _safe_data_path(version.stored_path)
    if not path.is_file():
        raise AppError("REPORT_OUTPUT_NOT_FOUND", "PPTX 样板文件已丢失", status_code=404)
    return run, version, path


def _validated_sections(sections: list[str]) -> list[str]:
    selected = list(dict.fromkeys(sections or DEFAULT_REPORT_SECTIONS))
    unknown = set(selected) - set(DEFAULT_REPORT_SECTIONS)
    if unknown:
        invalid = ", ".join(sorted(unknown))
        raise AppError("REPORT_SECTION_INVALID", f"存在不支持的报表区域：{invalid}")
    return selected


def _share_code(settings_payload: dict) -> str | None:
    value = settings_payload.get("share_product_code")
    return value.strip() if isinstance(value, str) and value.strip() else None


def _run_item(run: ReportRun, product_name: str, session: Session | None = None) -> ReportRunItem:
    current_version = (
        session.get(ReportFileVersion, run.current_version_id)
        if session is not None and run.current_version_id is not None
        else None
    )
    return ReportRunItem(
        id=run.id,
        definition_id=run.definition_id,
        fund_product_id=run.fund_product_id,
        product_name=product_name,
        template_key=run.template_key,
        report_date=run.report_date,
        status=run.status,
        output_filename=run.output_filename,
        current_version_id=run.current_version_id,
        current_version=current_version.version if current_version else None,
        template_version_id=run.template_version_id,
        error_stage=run.error_stage,
        error_code=run.error_code,
        error_message=run.error_message,
        create_time=run.create_time,
    )


def _batch_view(batch: ReportBatch) -> ReportBatchView:
    return ReportBatchView.model_validate(batch, from_attributes=True)


def _batch_item_view(item: ReportBatchItem, product_name: str) -> ReportBatchItemView:
    return ReportBatchItemView(
        id=item.id,
        fund_product_id=item.fund_product_id,
        product_name=product_name,
        status=item.status,
        report_run_id=item.report_run_id,
        attempt_count=item.attempt_count,
        error_code=item.error_code,
        error_message=item.error_message,
    )


def _audit(
    session: Session,
    request: Request,
    scope: TenantContext,
    *,
    action: str,
    resource_type: str,
    resource_id: int,
    detail: dict,
    outcome: str = "success",
) -> None:
    AuditService(audit_signing_key(get_settings().security)).append(
        session,
        tenant_id=scope.tenant_id,
        actor_user_id=scope.user.id,
        actor_username=scope.user.username,
        action=action,
        resource_type=resource_type,
        resource_id=resource_id,
        outcome=outcome,
        detail=detail,
        ip_address=request.client.host if request.client else None,
        user_agent=request.headers.get("user-agent"),
    )
