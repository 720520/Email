from __future__ import annotations

import base64
from copy import deepcopy
from datetime import date, timedelta
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.services.report_presentation_service import ReportPresentationService
from app.services.report_template_service import (
    ReportTemplateInspector,
    format_template_value,
)


def _snapshot() -> dict:
    start = date(2026, 1, 1)
    return {
        "product_name": "动态模板测试基金",
        "product_code": "TPL001",
        "report_date": "2026-08-21",
        "fields": {
            "product_name": "动态模板测试基金",
            "product_code": "TPL001",
            "investment_manager": "张经理",
            "manager_name": "测试管理机构",
        },
        "performance": {
            "annualized_return": "18.126%",
            "return_ytd": "12.50%",
        },
        "nav_series": [
            {
                "date": (start + timedelta(days=index * 30)).isoformat(),
                "unit_nav": str(1 + index / 100),
                "total_nav": str(1 + index / 100),
            }
            for index in range(8)
        ],
        "dynamic_fields": {"custom.roadshow_contact": "李经理"},
    }


def _template() -> Presentation:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.6))
    paragraph = title.text_frame.paragraphs[0]
    paragraph.add_run().text = "产品：{{product."
    paragraph.add_run().text = "name}} / {{product.name}}"
    contact = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(0.5))
    contact.text = '联系人：{{custom.roadshow_contact|default:"—"}}'
    metric = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(0.5))
    metric.text = "年化：{{metric.annualized_return|percent:2}}"
    table = slide.shapes.add_table(2, 2, Inches(0.5), Inches(2.4), Inches(8), Inches(1)).table
    table.cell(0, 0).text = "名称"
    table.cell(0, 1).text = "代码"
    table.cell(1, 0).text = "{{product.name}}"
    table.cell(1, 1).text = "{{product.code}}"
    chart_anchor = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(8), Inches(2.5))
    chart_anchor.text = "{{chart:nav_history}}"
    return presentation


def test_inspector_finds_repeated_split_run_table_and_component_tokens() -> None:
    inspection = ReportTemplateInspector().inspect(_template())

    assert inspection.is_valid
    assert inspection.required_fields == (
        "custom.roadshow_contact",
        "metric.annualized_return",
        "product.code",
        "product.name",
    )
    assert inspection.required_components == ("chart:nav_history",)
    assert sum(item.field_key == "product.name" for item in inspection.tokens) == 3


def test_inspector_rejects_unknown_formatter_and_duplicate_component() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1))
    first.text = "{{product.name|unsafe}}"
    second = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(2), Inches(1))
    second.text = "{{chart:nav_history}}"
    third = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(2), Inches(1))
    third.text = "{{chart:nav_history}}"

    inspection = ReportTemplateInspector().inspect(presentation)

    assert {item["code"] for item in inspection.errors} == {
        "UNKNOWN_FORMATTER",
        "DUPLICATE_COMPONENT",
    }


def test_inspector_scans_multiple_slides_and_layout_placeholders() -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    first.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1)).text = "{{product.name}}"
    second.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1)).text = "{{product.name}}"
    presentation.slide_layouts[0].placeholders[0].text = "{{report.date}}"

    inspection = ReportTemplateInspector().inspect(presentation)

    assert inspection.required_fields == ("product.name", "report.date")
    assert sum(item.field_key == "product.name" for item in inspection.tokens) == 2
    assert any(item.location.startswith("layout:") for item in inspection.tokens)


def test_dynamic_template_render_replaces_all_occurrences_and_chart(tmp_path: Path) -> None:
    template = tmp_path / "template.pptx"
    output = tmp_path / "output.pptx"
    _template().save(template)

    ReportPresentationService().generate(
        _snapshot(),
        output_path=output,
        sections=[],
        template_path=template,
    )

    rendered = Presentation(output)
    text = "\n".join(
        shape.text
        for slide in rendered.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    table_text = "|".join(
        cell.text
        for slide in rendered.slides
        for shape in slide.shapes
        if getattr(shape, "has_table", False)
        for row in shape.table.rows
        for cell in row.cells
    )
    charts = [
        shape
        for slide in rendered.slides
        for shape in slide.shapes
        if getattr(shape, "has_chart", False)
    ]
    assert "产品：动态模板测试基金 / 动态模板测试基金" in text
    assert "联系人：李经理" in text
    assert "年化：18.13%" in text
    assert "动态模板测试基金|TPL001" in table_text
    assert "{{" not in text + table_text
    assert len(charts) == 1


def test_two_products_render_independently_without_mutating_template(tmp_path: Path) -> None:
    template = tmp_path / "shared-template.pptx"
    first_output = tmp_path / "first.pptx"
    second_output = tmp_path / "second.pptx"
    _template().save(template)
    original = template.read_bytes()
    first_snapshot = _snapshot()
    second_snapshot = deepcopy(first_snapshot)
    second_snapshot["product_name"] = "第二只基金"
    second_snapshot["product_code"] = "TPL002"
    second_snapshot["fields"]["product_name"] = "第二只基金"
    second_snapshot["fields"]["product_code"] = "TPL002"
    second_snapshot["dynamic_fields"]["custom.roadshow_contact"] = "王经理"

    service = ReportPresentationService()
    service.generate(
        first_snapshot,
        output_path=first_output,
        sections=[],
        template_path=template,
    )
    service.generate(
        second_snapshot,
        output_path=second_output,
        sections=[],
        template_path=template,
    )

    def all_text(path: Path) -> str:
        presentation = Presentation(path)
        return "|".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )

    assert "动态模板测试基金" in all_text(first_output)
    assert "第二只基金" not in all_text(first_output)
    assert "第二只基金" in all_text(second_output)
    assert "动态模板测试基金" not in all_text(second_output)
    assert template.read_bytes() == original


def test_format_template_value_supports_date_default_and_percentage() -> None:
    assert format_template_value("2026-08-21", ("date",)) == "2026年08月21日"
    assert format_template_value(None, ('default:"未提供"',)) == "未提供"
    assert format_template_value("12.345%", ("percent:1",)) == "12.3%"


def test_image_anchor_and_empty_default_are_rendered(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(4), Inches(0.5)
    ).text = '{{custom.optional_note|default:"暂无"}}'
    slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(2), Inches(1)
    ).text = "{{image:custom.company_logo}}"
    template = tmp_path / "image-template.pptx"
    output = tmp_path / "image-output.pptx"
    presentation.save(template)
    snapshot = _snapshot()
    snapshot["dynamic_fields"].update(
        {
            "custom.optional_note": None,
            "custom.company_logo": "data:image/png;base64,"
            + base64.b64encode(
                base64.b64decode(
                    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
                )
            ).decode(),
        }
    )

    ReportPresentationService().generate(
        snapshot,
        output_path=output,
        sections=[],
        template_path=template,
    )

    rendered = Presentation(output)
    assert any(
        shape.text == "暂无"
        for shape in rendered.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert any(shape.shape_type == 13 for shape in rendered.slides[0].shapes)
