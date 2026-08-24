from __future__ import annotations

from fastapi import FastAPI
from pptx import Presentation
from pptx.util import Inches

from app.services.report_template_service import ReportTemplateInspector


def test_template_validation_uses_tenant_field_registry(app: FastAPI) -> None:
    del app
    from app.core.config import get_settings
    from app.db.models import ReportFieldDefinition
    from app.db.session import configure_tenant_scope, get_database_manager
    from app.services.foundation_service import FoundationService

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(8), Inches(1))
    box.text = "{{custom.known_field}} / {{custom.unknown_field}}"

    with get_database_manager().session_factory() as session, session.begin():
        identity = FoundationService(get_settings()).ensure(session)
        configure_tenant_scope(
            session,
            tenant_id=identity.tenant_id,
            mailbox_ids=(identity.mailbox_account_id,),
        )
        session.add(
            ReportFieldDefinition(
                tenant_id=identity.tenant_id,
                field_key="custom.known_field",
                label="已知字段",
                data_type="string",
                value_kind="scalar",
                source_type="custom",
            )
        )
        session.flush()
        inspection = ReportTemplateInspector().inspect(presentation, session)

    assert inspection.required_fields == (
        "custom.known_field",
        "custom.unknown_field",
    )
    assert [item["field_key"] for item in inspection.errors] == ["custom.unknown_field"]


def test_template_validation_checks_required_defaults_and_formatter_types(
    app: FastAPI,
) -> None:
    del app
    from app.core.config import get_settings
    from app.db.models import ReportFieldDefinition
    from app.db.session import configure_tenant_scope, get_database_manager
    from app.services.foundation_service import FoundationService

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(8), Inches(1))
    box.text = "{{custom.required_text|percent:2}} / {{custom.report_day|date}}"

    with get_database_manager().session_factory() as session, session.begin():
        identity = FoundationService(get_settings()).ensure(session)
        configure_tenant_scope(
            session,
            tenant_id=identity.tenant_id,
            mailbox_ids=(identity.mailbox_account_id,),
        )
        session.add_all(
            [
                ReportFieldDefinition(
                    tenant_id=identity.tenant_id,
                    field_key="custom.required_text",
                    label="必填文本",
                    data_type="string",
                    value_kind="scalar",
                    source_type="custom",
                    is_required=True,
                ),
                ReportFieldDefinition(
                    tenant_id=identity.tenant_id,
                    field_key="custom.report_day",
                    label="报告日期",
                    data_type="date",
                    value_kind="scalar",
                    source_type="custom",
                    default_value="2026-08-21",
                ),
            ]
        )
        session.flush()
        inspection = ReportTemplateInspector().inspect(presentation, session)

    assert {item["code"] for item in inspection.errors} == {
        "FORMATTER_TYPE_MISMATCH",
        "REQUIRED_FIELD_DEFAULT_MISSING",
    }
