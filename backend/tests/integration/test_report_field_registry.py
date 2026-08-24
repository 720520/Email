from __future__ import annotations

from datetime import date

import pytest
from fastapi import FastAPI
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy import select

from app.core.errors import AppError
from app.db.models import FundProduct, ReportFieldDefinition, ReportFieldValue, Tenant
from app.services.report_field_service import FieldContext, ReportFieldResolver
from app.services.report_template_service import ReportTemplateInspector


def test_custom_field_persistence_and_resolution(app: FastAPI) -> None:
    del app
    from app.core.config import get_settings
    from app.db.session import configure_tenant_scope, get_database_manager
    from app.services.foundation_service import FoundationService

    resolver = ReportFieldResolver()
    with get_database_manager().session_factory() as session, session.begin():
        identity = FoundationService(get_settings()).ensure(session)
        configure_tenant_scope(
            session,
            tenant_id=identity.tenant_id,
            mailbox_ids=(identity.mailbox_account_id,),
        )
        product = FundProduct(
            tenant_id=identity.tenant_id,
            product_code="REGISTRY001",
            product_name="字段注册中心测试基金",
            source_profile={"manager_name": "测试管理人"},
            source_profile_meta={
                "manager_name": {
                    "source_type": "contract",
                    "source_reference": "contract.pdf",
                }
            },
        )
        definition = ReportFieldDefinition(
            tenant_id=identity.tenant_id,
            field_key="custom.roadshow_contact",
            label="路演联系人",
            data_type="string",
            value_kind="scalar",
            source_type="custom",
            default_value="未配置",
        )
        session.add_all([product, definition])
        session.flush()
        session.add(
            ReportFieldValue(
                tenant_id=identity.tenant_id,
                field_definition_id=definition.id,
                entity_type="fund_product",
                entity_id=product.id,
                value_text="张经理",
                effective_date=date(2026, 8, 1),
                source_type="manual",
                source_reference="路演资料",
            )
        )
        session.flush()

        resolved = resolver.resolve_many(
            session,
            ["product.name", "product.manager_name", "custom.roadshow_contact"],
            FieldContext(
                tenant_id=identity.tenant_id,
                tenant_name="默认业务账套",
                product_id=product.id,
                report_date=date(2026, 8, 21),
            ),
        )

        assert resolved["product.name"][1].value == "字段注册中心测试基金"
        assert resolved["product.manager_name"][1].source_reference == "contract.pdf"
        assert resolved["custom.roadshow_contact"][1].value == "张经理"
        assert resolved["custom.roadshow_contact"][1].source_reference == "路演资料"

        definition.is_active = False
        session.flush()
        with pytest.raises(AppError) as error:
            resolver.resolve_many(
                session,
                ["custom.roadshow_contact"],
                FieldContext(identity.tenant_id, "默认业务账套", product.id),
            )
        assert error.value.code == "REPORT_FIELD_NOT_FOUND"
        historical = resolver.resolve_many(
            session,
            ["custom.roadshow_contact"],
            FieldContext(identity.tenant_id, "默认业务账套", product.id),
            allow_inactive=True,
        )
        assert historical["custom.roadshow_contact"][1].value == "张经理"


def test_custom_field_uses_default_when_product_value_missing(app: FastAPI) -> None:
    del app
    from app.core.config import get_settings
    from app.db.session import configure_tenant_scope, get_database_manager
    from app.services.foundation_service import FoundationService

    with get_database_manager().session_factory() as session, session.begin():
        identity = FoundationService(get_settings()).ensure(session)
        configure_tenant_scope(
            session,
            tenant_id=identity.tenant_id,
            mailbox_ids=(identity.mailbox_account_id,),
        )
        product = FundProduct(
            tenant_id=identity.tenant_id,
            product_code="DEFAULT001",
            product_name="默认值测试基金",
        )
        session.add(product)
        session.add(
            ReportFieldDefinition(
                tenant_id=identity.tenant_id,
                field_key="custom.disclaimer_version",
                label="免责声明版本",
                data_type="string",
                value_kind="scalar",
                source_type="custom",
                default_value="v1",
            )
        )
        session.flush()
        resolved = ReportFieldResolver().resolve_many(
            session,
            ["custom.disclaimer_version"],
            FieldContext(identity.tenant_id, "默认业务账套", product.id),
        )["custom.disclaimer_version"][1]
        assert resolved.value == "v1"
        assert resolved.used_default is True


def test_custom_fields_are_isolated_between_tenants(app: FastAPI) -> None:
    del app
    from app.core.config import get_settings
    from app.db.session import configure_tenant_scope, get_database_manager
    from app.services.foundation_service import FoundationService

    with get_database_manager().session_factory() as session:
        identity = FoundationService(get_settings()).ensure(session)
        session.info["skip_tenant_scope"] = True
        other_tenant = Tenant(code="other-field-tenant", name="其他字段租户")
        session.add(other_tenant)
        session.flush()
        other_field = ReportFieldDefinition(
            tenant_id=other_tenant.id,
            field_key="custom.private_contact",
            label="其他租户联系人",
            data_type="string",
            value_kind="scalar",
            source_type="custom",
        )
        session.add(other_field)
        session.commit()
        other_field_id = other_field.id
        session.info.pop("skip_tenant_scope")
        configure_tenant_scope(
            session,
            tenant_id=identity.tenant_id,
            mailbox_ids=(identity.mailbox_account_id,),
        )

        with pytest.raises(AppError) as error:
            ReportFieldResolver().definition(session, "custom.private_contact")
        assert error.value.code == "REPORT_FIELD_NOT_FOUND"
        assert session.scalar(
            select(ReportFieldDefinition).where(ReportFieldDefinition.id == other_field_id)
        ) is None
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(4), Inches(1)
        ).text = "{{custom.private_contact}}"
        inspection = ReportTemplateInspector().inspect(presentation, session)
        assert inspection.errors[0]["code"] == "UNKNOWN_FIELD"
