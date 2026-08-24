from __future__ import annotations

import hashlib
import io
import tempfile
from datetime import date
from decimal import Decimal

import pytest
from fastapi import FastAPI, UploadFile
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy import select
from starlette.requests import Request

from app.api.deps import TenantContext
from app.api.schemas.reporting import ReportGenerateRequest
from app.api.v1.report_fields import field_usages
from app.api.v1.reports import (
    download_report,
    generate_report,
    list_report_templates,
    publish_report_template,
    regenerate_report_from_snapshot,
    upload_report_template,
    upload_report_template_version,
)
from app.core.errors import AppError
from app.db.models import (
    AuditEvent,
    FundNav,
    FundProduct,
    ReportFieldDefinition,
    ReportFieldValue,
    ReportFileVersion,
    UserRole,
)


def _request(path: str) -> Request:
    return Request(
        {
            "type": "http",
            "method": "POST",
            "path": path,
            "headers": [],
            "client": ("127.0.0.1", 12345),
        }
    )


def _template_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    first.text = "{{product.name}} / {{product.name}}"
    second = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1))
    second.text = "路演联系人：{{custom.roadshow_contact}}"
    content = io.BytesIO()
    presentation.save(content)
    return content.getvalue()


def _template_upload() -> UploadFile:
    file = tempfile.SpooledTemporaryFile(max_size=2 * 1024 * 1024)
    file.write(_template_bytes())
    file.seek(0)
    return UploadFile(filename="dynamic.pptx", file=file)


def _upload_bytes(content: bytes, filename: str = "dynamic.pptx") -> UploadFile:
    file = tempfile.SpooledTemporaryFile(max_size=2 * 1024 * 1024)
    file.write(content)
    file.seek(0)
    return UploadFile(filename=filename, file=file)


def _unknown_field_template() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(8), Inches(1)
    ).text = "{{custom.not_registered}}"
    content = io.BytesIO()
    presentation.save(content)
    return content.getvalue()


@pytest.mark.anyio
async def test_template_draft_publish_and_dynamic_generation_loop(
    app: FastAPI, monkeypatch: pytest.MonkeyPatch
) -> None:
    del app
    from app.core.config import get_settings
    from app.db.session import configure_tenant_scope, get_database_manager
    from app.services.auth_service import AuthService
    from app.services.foundation_service import FoundationService

    with get_database_manager().session_factory() as session:
        identity = FoundationService(get_settings()).ensure(session)
        user = AuthService().create_user(
            session,
            username="template-admin",
            password="AdminPass!2026",
            role=UserRole.ADMIN,
            tenant_id=identity.tenant_id,
        )
        configure_tenant_scope(
            session,
            tenant_id=identity.tenant_id,
            mailbox_ids=(identity.mailbox_account_id,),
        )
        product = FundProduct(
            tenant_id=identity.tenant_id,
            product_code="TEMPLATE001",
            product_name="模板生命周期测试基金",
        )
        field = ReportFieldDefinition(
            tenant_id=identity.tenant_id,
            field_key="custom.roadshow_contact",
            label="路演联系人",
            data_type="string",
            value_kind="scalar",
            source_type="custom",
        )
        session.add_all([product, field])
        session.flush()
        session.add_all(
            [
                FundNav(
                    tenant_id=identity.tenant_id,
                    mailbox_account_id=identity.mailbox_account_id,
                    product_name=product.product_name,
                    product_code=product.product_code,
                    master_product_code=product.product_code,
                    nav_date=date(2026, 8, 21),
                    unit_nav=Decimal("1.1000"),
                    total_nav=Decimal("1.1000"),
                    source_file="nav.xlsx",
                ),
                ReportFieldValue(
                    tenant_id=identity.tenant_id,
                    field_definition_id=field.id,
                    entity_type="fund_product",
                    entity_id=product.id,
                    value_text="王经理",
                    source_type="manual",
                ),
            ]
        )
        session.commit()
        scope = TenantContext(
            user=user,
            tenant_id=identity.tenant_id,
            tenant_code="default",
            tenant_name="默认业务账套",
            role=UserRole.ADMIN,
            mailbox_ids=(identity.mailbox_account_id,),
            content_mailbox_ids=(identity.mailbox_account_id,),
            operable_mailbox_ids=(identity.mailbox_account_id,),
            manageable_mailbox_ids=(identity.mailbox_account_id,),
        )

        draft = await upload_report_template(
            request=_request("/api/v1/reports/templates"),
            session=session,
            scope=scope,
            file=_template_upload(),
            name="动态字段模板",
            description="测试草稿发布闭环",
        )
        assert draft.status == "draft"
        assert draft.validation_errors == []
        assert draft.required_fields == ["custom.roadshow_contact", "product.name"]

        published = publish_report_template(
            draft.id or 0,
            _request(f"/api/v1/reports/templates/{draft.id}/publish"),
            session,
            scope,
        )
        generated = generate_report(
            ReportGenerateRequest(
                fund_product_id=product.id,
                template_key=published.key,
                report_date=date(2026, 8, 21),
                sections=[],
            ),
            _request("/api/v1/reports/generate"),
            session,
            scope,
        )

        # API 返回文件名，真实存储路径由 ReportRun 保存。
        from app.db.models import ReportRun

        run = session.get(ReportRun, generated.run.id)
        assert run is not None and run.output_path
        output_path = get_settings().data_directory / run.output_path
        rendered = Presentation(output_path)
        text = "\n".join(
            shape.text
            for slide in rendered.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert published.status == "published"
        assert published.key.startswith("template-version:")
        assert text.count("模板生命周期测试基金") == 2
        assert "路演联系人：王经理" in text
        assert "{{" not in text
        original_snapshot = dict(run.input_snapshot)
        product.product_name = "数据库已修改的基金名称"
        contact_value = session.scalar(
            select(ReportFieldValue).where(
                ReportFieldValue.field_definition_id == field.id,
                ReportFieldValue.entity_id == product.id,
            )
        )
        assert contact_value is not None
        contact_value.value_text = "数据库已修改联系人"
        session.commit()
        regenerated = regenerate_report_from_snapshot(
            run.id,
            _request(f"/api/v1/reports/runs/{run.id}/regenerate"),
            session,
            scope,
        )
        versions = list(
            session.scalars(
                select(ReportFileVersion)
                .where(ReportFileVersion.report_run_id == run.id)
                .order_by(ReportFileVersion.version)
            )
        )
        latest = Presentation(get_settings().data_directory / versions[-1].stored_path)
        latest_text = "\n".join(
            shape.text
            for slide in latest.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert regenerated.run.current_version == 2
        assert [version.version for version in versions] == [1, 2]
        assert all(len(version.content_hash) == 64 for version in versions)
        assert (
            versions[-1].content_hash
            == hashlib.sha256(
                (get_settings().data_directory / versions[-1].stored_path).read_bytes()
            ).hexdigest()
        )
        assert run.input_snapshot == original_snapshot
        assert "模板生命周期测试基金" in latest_text
        assert "数据库已修改的基金名称" not in latest_text
        assert "王经理" in latest_text
        assert "数据库已修改联系人" not in latest_text
        response = download_report(
            run.id,
            _request(f"/api/v1/reports/runs/{run.id}/download"),
            session,
            scope,
        )
        assert response.path == get_settings().data_directory / versions[-1].stored_path
        assert (
            session.scalar(
                select(AuditEvent).where(
                    AuditEvent.action == "report.download",
                    AuditEvent.resource_id == versions[-1].id,
                )
            )
            is not None
        )
        usages = field_usages("custom.roadshow_contact", session, scope)
        assert usages["template_count"] == 1
        assert usages["templates"][0]["version"] == 1
        with pytest.raises(AppError) as immutable_error:
            publish_report_template(
                draft.id or 0,
                _request(f"/api/v1/reports/templates/{draft.id}/publish"),
                session,
                scope,
            )
        assert immutable_error.value.code == "REPORT_TEMPLATE_DRAFT_NOT_FOUND"

        next_draft = await upload_report_template_version(
            draft.id or 0,
            _request(f"/api/v1/reports/templates/{draft.id}/versions"),
            session,
            scope,
            _template_upload(),
        )
        visible_versions = [
            item for item in list_report_templates(session, scope) if item.id == draft.id
        ]
        assert next_draft.version == 2
        assert next_draft.status == "draft"
        assert [(item.version, item.status) for item in visible_versions] == [
            (1, "published"),
            (2, "draft"),
        ]

        invalid_draft = await upload_report_template(
            request=_request("/api/v1/reports/templates"),
            session=session,
            scope=scope,
            file=_upload_bytes(_unknown_field_template(), "unknown-field.pptx"),
            name="无效字段模板",
            description=None,
        )
        assert invalid_draft.validation_errors[0]["code"] == "UNKNOWN_FIELD"
        with pytest.raises(AppError) as publish_error:
            publish_report_template(
                invalid_draft.id or 0,
                _request(f"/api/v1/reports/templates/{invalid_draft.id}/publish"),
                session,
                scope,
            )
        assert publish_error.value.code == "REPORT_TEMPLATE_VALIDATION_FAILED"

        with pytest.raises(AppError) as corrupt_error:
            await upload_report_template(
                request=_request("/api/v1/reports/templates"),
                session=session,
                scope=scope,
                file=_upload_bytes(b"not-a-pptx", "broken.pptx"),
                name="损坏模板",
                description=None,
            )
        assert corrupt_error.value.code == "REPORT_TEMPLATE_INVALID"

        def fail_render(*args, **kwargs) -> None:
            del args, kwargs
            raise RuntimeError("模拟渲染失败")

        monkeypatch.setattr("app.api.v1.reports.ReportPresentationService.generate", fail_render)
        with pytest.raises(AppError) as generation_error:
            generate_report(
                ReportGenerateRequest(
                    fund_product_id=product.id,
                    template_key=published.key,
                    report_date=date(2026, 8, 21),
                    sections=[],
                ),
                _request("/api/v1/reports/generate"),
                session,
                scope,
            )
        failed_run = session.scalar(select(ReportRun).order_by(ReportRun.id.desc()))
        assert generation_error.value.code == "REPORT_GENERATION_FAILED"
        assert failed_run is not None
        assert failed_run.error_stage == "render"
        assert failed_run.error_code == "REPORT_RENDER_FAILED"
        assert failed_run.input_snapshot
        assert not list(get_settings().data_directory.rglob("*.tmp.pptx"))
